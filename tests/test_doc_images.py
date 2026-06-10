"""文書内画像の抽出・保存・チャンク紐付け(doc_images)と図チャンク生成のテスト。

Excel/Word/PowerPoint は実ファイルを生成して抽出を検証する。PDF は pymupdf が
ある環境のみ(無ければ自動スキップ)。Ollama・埋め込みモデルには接続しない
(埋め込みと図の説明はフェイクに差し替え、ChromaDB は一時フォルダで実体を使う)。
pytest でも `python tests/test_doc_images.py` でも動く。
"""
import contextlib
import io
import json
import os
import shutil
import sys
import tempfile
from pathlib import Path

sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

from app import db, doc_images, ocr, rag       # noqa: E402
from app.config import settings                 # noqa: E402


def _has(mod: str) -> bool:
    try:
        __import__(mod)
        return True
    except Exception:
        return False


def _png_bytes(w: int = 320, h: int = 200) -> bytes:
    """乱数ノイズ画像(単色・規則模様だとPNG圧縮で MIN_BYTES 未満になり実画像の代わりにならない)。"""
    import random
    from PIL import Image
    raw = random.Random(42).randbytes(w * h * 3)
    buf = io.BytesIO()
    Image.frombytes("RGB", (w, h), raw).save(buf, format="PNG")
    return buf.getvalue()


@contextlib.contextmanager
def patched(obj, name, value):
    old = getattr(obj, name)
    setattr(obj, name, value)
    try:
        yield
    finally:
        setattr(obj, name, old)


@contextlib.contextmanager
def temp_env():
    """data_dir / chroma_dir / db を一時フォルダへ差し替え、資料用フォルダも用意する。"""
    d = Path(tempfile.mkdtemp(prefix="docimg_test_")).resolve()
    (d / "data").mkdir()
    ws = d / "docs"
    ws.mkdir()
    old = (settings.data_dir, settings.chroma_dir, settings.db_path)
    settings.data_dir = d / "data"
    settings.chroma_dir = d / "data" / "chroma"
    settings.db_path = d / "data" / "t.db"
    old_client = rag._client
    rag._client = None
    try:
        db.init_db()
        yield ws
    finally:
        rag._client = old_client
        settings.data_dir, settings.chroma_dir, settings.db_path = old
        shutil.rmtree(d, ignore_errors=True)


class FakeEmb:
    def embed_query(self, t):
        return [0.1] * 8

    def embed_documents(self, ts):
        return [[0.1] * 8 for _ in ts]


# ---------------- 形式判定・正規化 ----------------
def test_sniff_ext():
    assert doc_images._sniff_ext(_png_bytes()) == "png"
    assert doc_images._sniff_ext(b"\xff\xd8\xff\xe0" + b"x" * 10) == "jpg"
    assert doc_images._sniff_ext(b"GIF89a" + b"x" * 10) == "gif"
    assert doc_images._sniff_ext(b"not an image") == ""


def test_normalize_filters_tiny_and_resizes_huge():
    assert doc_images._normalize(_png_bytes(16, 16)) is None        # アイコンは除外
    assert doc_images._normalize(b"x" * 100) is None                # 小さすぎ・非画像
    big = doc_images._normalize(_png_bytes(2400, 1200))
    assert big is not None
    from PIL import Image
    img = Image.open(io.BytesIO(big[0]))
    assert max(img.size) <= doc_images.MAX_DIM                      # 縮小される
    ok = doc_images._normalize(_png_bytes(320, 200))
    assert ok is not None and ok[1] == "png"                        # 通常サイズは素通し


# ---------------- Excel ----------------
def _make_xlsx(path: Path, png_path: Path, second_sheet: bool = False):
    import openpyxl
    from openpyxl.drawing.image import Image as XLImage
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "組立"
    ws["A1"] = "手順"
    ws["A2"] = "図1のとおりネジを締める"
    ws.add_image(XLImage(str(png_path)), "B2")
    if second_sheet:
        ws2 = wb.create_sheet("配線")
        ws2["A1"] = "図2参照"
        ws2.add_image(XLImage(str(png_path)), "A3")
    wb.save(str(path))


def test_xlsx_extract_links_to_sheet():
    with temp_env() as ws:
        png = ws / "fig.png"
        png.write_bytes(_png_bytes())
        _make_xlsx(ws / "手順書.xlsx", png)
        loc_map, images = doc_images.extract_for_file("idx1", ws / "手順書.xlsx")
        assert "シート:組立" in loc_map and len(loc_map["シート:組立"]) == 1
        img_id = loc_map["シート:組立"][0]
        assert img_id.startswith("idx1/")
        iid, name = img_id.split("/", 1)
        assert (settings.data_dir / "doc_images" / iid / name).is_file()
        assert images and images[0][0] == img_id and images[0][2] == "シート:組立"


def test_xlsx_same_image_two_sheets_saved_once():
    with temp_env() as ws:
        png = ws / "fig.png"
        png.write_bytes(_png_bytes())
        _make_xlsx(ws / "手順書.xlsx", png, second_sheet=True)
        loc_map, images = doc_images.extract_for_file("idx2", ws / "手順書.xlsx")
        assert set(loc_map) == {"シート:組立", "シート:配線"}
        assert loc_map["シート:組立"] == loc_map["シート:配線"]   # 同一内容 → 同一ID
        files = list((settings.data_dir / "doc_images" / "idx2").glob("*"))
        assert len(files) == 1                                      # 物理ファイルは1つ
        assert len(images) == 1                                     # 説明生成も1回分


# ---------------- Word ----------------
def test_docx_extract_associates_heading():
    with temp_env() as ws:
        from docx import Document
        png = ws / "fig.png"
        png.write_bytes(_png_bytes())
        d = Document()
        d.add_paragraph("前置きの文章")
        d.add_heading("設置手順", level=1)
        d.add_picture(str(png))
        d.save(str(ws / "マニュアル.docx"))
        loc_map, images = doc_images.extract_for_file("idx3", ws / "マニュアル.docx")
        assert list(loc_map) == ["設置手順"]                       # 直前の見出しに紐づく
        assert len(images) == 1


# ---------------- PowerPoint ----------------
def test_pptx_extract_links_to_slide():
    with temp_env() as ws:
        from pptx import Presentation
        from pptx.util import Inches
        png = ws / "fig.png"
        png.write_bytes(_png_bytes())
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.shapes.add_picture(str(png), Inches(1), Inches(1))
        prs.save(str(ws / "説明資料.pptx"))
        loc_map, _ = doc_images.extract_for_file("idx4", ws / "説明資料.pptx")
        assert "スライド1" in loc_map


# ---------------- PDF(pymupdf がある環境のみ)----------------
def test_pdf_extract_links_to_page():
    if not _has("fitz"):
        return   # pymupdf 無し環境は自動スキップ
    with temp_env() as ws:
        import fitz
        png = ws / "fig.png"
        png.write_bytes(_png_bytes())
        pdf = fitz.open()
        page = pdf.new_page()
        page.insert_text((72, 72), "図1参照")
        page.insert_image(fitz.Rect(72, 100, 372, 300), filename=str(png))
        pdf.save(str(ws / "図面.pdf"))
        pdf.close()
        loc_map, _ = doc_images.extract_for_file("idx5", ws / "図面.pdf")
        assert "p.1" in loc_map and len(loc_map["p.1"]) == 1


# ---------------- build_index 統合(リンク + 図チャンク + 検索) ----------------
def test_build_index_links_images_and_creates_fig_chunks():
    with temp_env() as ws:
        png = ws / "fig.png"
        png.write_bytes(_png_bytes())
        _make_xlsx(ws / "手順書.xlsx", png)
        (ws / "memo.txt").write_text("ただのメモです", encoding="utf-8")
        idx = db.create_index("資料", [str(ws)])
        with patched(rag, "get_embedder", lambda: FakeEmb()), \
             patched(ocr, "describe_image_png", lambda b: "ネジ締結手順を示す組立図。ラベル: 図1"):
            row = rag.build_index(idx["id"], [str(ws)])
        assert row["status"] == "ready", row.get("error")

        col = rag._collection(f"idx_{idx['id']}")
        got = col.get(include=["metadatas", "documents"])
        # 本文チャンクに画像が紐づく
        linked = [m for m in got["metadatas"] if m.get("images")]
        assert linked, "画像が紐づいたチャンクがあるはず"
        ids = json.loads(linked[0]["images"])
        iid, name = ids[0].split("/", 1)
        assert (settings.data_dir / "doc_images" / iid / name).is_file()
        # 図チャンク(VLM説明)が索引化される
        figs = [d for d in got["documents"] if d.startswith("〔図〕")]
        assert figs and "ネジ締結" in figs[0]
        # 検索ヒットに images が載って出典へ流れる
        with patched(rag, "get_embedder", lambda: FakeEmb()):
            hits = rag.retrieve("ネジ締結の図", [idx["id"]], None, 5)
        assert any(h.get("images") for h in hits)

        # インデックス削除で画像フォルダも消える
        from app.services import index_service
        assert index_service.delete_index(idx["id"]) is True
        assert not (settings.data_dir / "doc_images" / idx["id"]).exists()


def test_build_index_without_vlm_still_links_images():
    """図の説明が空(OCR無効・非対応モデル)でも、画像リンクだけは付く。"""
    with temp_env() as ws:
        png = ws / "fig.png"
        png.write_bytes(_png_bytes())
        _make_xlsx(ws / "手順書.xlsx", png)
        idx = db.create_index("資料", [str(ws)])
        with patched(rag, "get_embedder", lambda: FakeEmb()), \
             patched(ocr, "describe_image_png", lambda b: ""):
            row = rag.build_index(idx["id"], [str(ws)])
        assert row["status"] == "ready"
        got = rag._collection(f"idx_{idx['id']}").get(include=["metadatas", "documents"])
        assert any(m.get("images") for m in got["metadatas"])
        assert not any(d.startswith("〔図〕") for d in got["documents"])


# ---------------- 配信ルートの安全性 ----------------
def test_doc_image_route_blocks_traversal_and_serves_valid():
    from fastapi import HTTPException
    from app.routes.index_routes import api_doc_image
    with temp_env():
        data = _png_bytes()
        img_id = doc_images._save("abcd1234", data, "png")
        iid, name = img_id.split("/", 1)
        # 正常系
        resp = api_doc_image(iid, name)
        assert Path(resp.path).is_file()
        # パストラバーサル・不正名は404
        for bad_iid, bad_name in ((".." , name), (iid, "../secret.png"),
                                  (iid, "evil.html"), ("x/y", name)):
            try:
                api_doc_image(bad_iid, bad_name)
                assert False, f"404になるはず: {bad_iid}/{bad_name}"
            except HTTPException as e:
                assert e.status_code == 404


if __name__ == "__main__":
    fns = [v for k, v in sorted(globals().items()) if k.startswith("test_") and callable(v)]
    passed = 0
    for fn in fns:
        try:
            fn()
            passed += 1
        except AssertionError as e:
            print(f"FAIL {fn.__name__}: {e}")
        except Exception as e:  # noqa: BLE001
            print(f"ERROR {fn.__name__}: {e!r}")
    print(f"{passed}/{len(fns)} passed")
    sys.exit(0 if passed == len(fns) else 1)
