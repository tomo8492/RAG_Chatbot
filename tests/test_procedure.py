"""手順ビューア(procedure.py)のテスト。

Excelの工程番号列の自動検出(1,2,3 / ①② / 偽陽性の排除)、画像のアンカー行による
工程への割り当て、フォールバック、PPT/PDF/Wordの単位分解、ルートのパス検証を固定する。
pytest でも `python tests/test_procedure.py` でも動く。
"""
import contextlib
import io
import os
import random
import shutil
import sys
import tempfile
from pathlib import Path

sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

from app import procedure                 # noqa: E402
from app.config import settings           # noqa: E402


def _has(mod: str) -> bool:
    try:
        __import__(mod)
        return True
    except Exception:
        return False


@contextlib.contextmanager
def temp_env():
    d = Path(tempfile.mkdtemp(prefix="proc_test_")).resolve()
    (d / "data").mkdir()
    docs = d / "docs"
    docs.mkdir()
    old = settings.data_dir
    settings.data_dir = d / "data"
    try:
        yield docs
    finally:
        settings.data_dir = old
        shutil.rmtree(d, ignore_errors=True)


def _png(path: Path, seed: int = 0):
    from PIL import Image
    raw = random.Random(seed).randbytes(200 * 140 * 3)
    buf = io.BytesIO()
    Image.frombytes("RGB", (200, 140), raw).save(buf, format="PNG")
    path.write_bytes(buf.getvalue())


def _make_xlsx(path: Path, png: Path, rows: dict, img_rows: list, sheet: str = "手順"):
    """rows={行番号(1始まり): [セル...]}, img_rows=[画像を貼る行(1始まり)]"""
    import openpyxl
    from openpyxl.drawing.image import Image as XLImage
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = sheet
    for r, vals in rows.items():
        for c, v in enumerate(vals, 1):
            ws.cell(row=r, column=c, value=v)
    for r in img_rows:
        ws.add_image(XLImage(str(png)), f"E{r}")
    wb.save(str(path))


# ---------------- Excel: 工程番号列の検出と画像の割り当て ----------------
def test_xlsx_steps_with_anchored_images():
    with temp_env() as docs:
        png = docs / "f.png"
        _png(png)
        _make_xlsx(docs / "手順書.xlsx", png,
                   rows={1: ["工程", "内容"],
                         2: ["1", "投入予定表を開く"],
                         4: ["", "補足: メニューから選択"],
                         5: ["2", "データを確認する"],
                         8: ["3", "保存する"]},
                   img_rows=[2, 6, 8])
        view = procedure.build_view("idxP", docs / "手順書.xlsx")
        sheet = view["sheets"][0]
        assert sheet["fallback"] is False
        steps = sheet["steps"]
        nos = [s["no"] for s in steps]
        assert nos == ["", "1", "2", "3"]                    # 前置き + 工程1〜3
        s1, s2, s3 = steps[1], steps[2], steps[3]
        assert s1["title"].startswith("投入予定表")
        assert "補足" in s1["text"]                          # 工程1の行範囲(2〜4行)を含む
        assert len(s1["images"]) == 1                        # 行2の画像
        assert len(s2["images"]) == 1                        # 行6の画像は工程2(5〜7行)
        assert len(s3["images"]) == 1                        # 行8の画像
        assert all(u.startswith("/api/doc-images/idxP/") for s in steps for u in s["images"])


def test_xlsx_circled_numbers():
    with temp_env() as docs:
        png = docs / "f.png"
        _png(png, seed=1)
        _make_xlsx(docs / "m.xlsx", png,
                   rows={1: ["①", "電源を入れる"], 3: ["②", "起動を確認"]},
                   img_rows=[1])
        sheet = procedure.build_view("idxC", docs / "m.xlsx")["sheets"][0]
        assert [s["no"] for s in sheet["steps"]] == ["1", "2"]
        assert len(sheet["steps"][0]["images"]) == 1


def test_xlsx_fallback_when_no_markers():
    with temp_env() as docs:
        png = docs / "f.png"
        _png(png, seed=2)
        _make_xlsx(docs / "memo.xlsx", png,
                   rows={1: ["これはただのメモ"], 2: ["番号列はない"]},
                   img_rows=[2])
        sheet = procedure.build_view("idxF", docs / "memo.xlsx")["sheets"][0]
        assert sheet["fallback"] is True
        assert len(sheet["steps"]) == 1 and len(sheet["steps"][0]["images"]) == 1


def test_xlsx_random_number_column_not_treated_as_steps():
    """数値データ列(昇順でない)を工程番号と誤検出しない。"""
    with temp_env() as docs:
        png = docs / "f.png"
        _png(png, seed=3)
        _make_xlsx(docs / "data.xlsx", png,
                   rows={1: ["100", "数量A"], 2: ["5", "数量B"],
                         3: ["87", "数量C"], 4: ["3", "数量D"]},
                   img_rows=[])
        sheet = procedure.build_view("idxR", docs / "data.xlsx")["sheets"][0]
        assert sheet["fallback"] is True


# ---------------- PowerPoint / Word / PDF ----------------
def test_pptx_slides_become_steps():
    with temp_env() as docs:
        from pptx import Presentation
        from pptx.util import Inches
        png = docs / "f.png"
        _png(png, seed=4)
        prs = Presentation()
        for i in range(2):
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            slide.shapes.add_picture(str(png), Inches(1), Inches(1))
        prs.save(str(docs / "s.pptx"))
        sheet = procedure.build_view("idxS", docs / "s.pptx")["sheets"][0]
        assert [s["no"] for s in sheet["steps"]] == ["1", "2"]
        assert all(len(s["images"]) == 1 for s in sheet["steps"])


def test_docx_headings_become_steps_with_table_image():
    with temp_env() as docs:
        from docx import Document
        png = docs / "f.png"
        _png(png, seed=5)
        d = Document()
        d.add_heading("設置手順", level=1)
        d.add_paragraph("床面を水平にする。")
        tbl = d.add_table(rows=1, cols=1)
        tbl.rows[0].cells[0].paragraphs[0].add_run().add_picture(str(png))
        d.add_heading("配線手順", level=1)
        d.add_paragraph("端子台に接続する。")
        d.save(str(docs / "m.docx"))
        sheet = procedure.build_view("idxW", docs / "m.docx")["sheets"][0]
        titles = [s["title"] for s in sheet["steps"]]
        assert titles == ["設置手順", "配線手順"]
        assert len(sheet["steps"][0]["images"]) == 1          # 表中の画像は直前の見出し節へ
        assert "床面" in sheet["steps"][0]["text"]


def test_pdf_pages_become_steps():
    if not _has("fitz"):
        return   # pymupdf 無し環境は自動スキップ
    with temp_env() as docs:
        import fitz
        png = docs / "f.png"
        _png(png, seed=6)
        pdf = fitz.open()
        for i in range(2):
            page = pdf.new_page()
            page.insert_text((72, 72), f"step {i + 1}")
            page.insert_image(fitz.Rect(72, 100, 272, 240), filename=str(png))
        pdf.save(str(docs / "p.pdf"))
        pdf.close()
        sheet = procedure.build_view("idxD", docs / "p.pdf")["sheets"][0]
        assert [s["no"] for s in sheet["steps"]] == ["1", "2"]
        assert all(s["images"] for s in sheet["steps"])


def test_unsupported_extension_raises():
    with temp_env() as docs:
        (docs / "a.txt").write_text("x", encoding="utf-8")
        try:
            procedure.build_view("x", docs / "a.txt")
            assert False, "ValueError になるはず"
        except ValueError as e:
            assert "未対応" in str(e)


# ---------------- ルートのパス検証(登録フォルダ外を拒否) ----------------
def test_resolve_index_file_rejects_outside_paths():
    from fastapi import HTTPException
    from app.routes.index_routes import _resolve_index_file
    with temp_env() as docs:
        inside = docs / "ok.xlsx"
        inside.write_bytes(b"PK")
        idx = {"paths": [str(docs)]}
        assert _resolve_index_file(idx, str(inside)) == inside.resolve()
        outside = Path(tempfile.gettempdir()) / "proc_outside.xlsx"
        outside.write_bytes(b"PK")
        try:
            try:
                _resolve_index_file(idx, str(outside))
                assert False, "登録フォルダ外は拒否されるはず"
            except HTTPException as e:
                assert e.status_code == 400
        finally:
            outside.unlink(missing_ok=True)


if __name__ == "__main__":
    fns = [v for k, v in sorted(globals().items()) if k.startswith("test_") and callable(v)]
    passed = 0
    for fn in fns:
        try:
            fn()
            passed += 1
        except AssertionError as e:
            print(f"FAIL {fn.__name__}: {e}")
        except Exception as e:  # noqa: BLE001
            print(f"ERROR {fn.__name__}: {e!r}")
    print(f"{passed}/{len(fns)} passed")
    sys.exit(0 if passed == len(fns) else 1)
