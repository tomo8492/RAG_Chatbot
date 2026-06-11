"""
export.py
回答(Markdown)を各種ファイル形式へ変換する。
  対応: md / txt / html / docx / xlsx / pptx / code(任意拡張子)
重い依存(python-docx等)は関数内で遅延インポート。
"""
from __future__ import annotations

import base64
import html as _html
import io
import re
import unicodedata
from datetime import datetime

from ..logging_setup import get_logger
from .blocks import (parse_blocks, _LINK, _align_style, _render_list, _item_plain, _item_level, _strip_inline, _inline_html, _inline_runs)

log = get_logger("export")

MIME = {
    "md": "text/markdown; charset=utf-8",
    "txt": "text/plain; charset=utf-8",
    "html": "text/html; charset=utf-8",
    "code": "text/plain; charset=utf-8",
    "pdf": "application/pdf",
    "csv": "text/csv; charset=utf-8",
    "docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    "xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    "pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
}
EXT = {"md": "md", "txt": "txt", "html": "html",
       "docx": "docx", "xlsx": "xlsx", "pptx": "pptx"}



# ============================================================
#  各形式レンダラ
# ============================================================
_HTML_CSS = """
:root{
  --ink:#1f2328; --muted:#5b6570; --accent:#3b5bdb; --accent-soft:rgba(59,91,219,.08);
  --line:#e7e9ee; --soft:#f6f8fb; --zebra:#fafbfd; --paper:#ffffff;
  --radius:10px; --maxw:820px;
  --sans:"Yu Gothic UI","Hiragino Kaku Gothic ProN","Hiragino Sans",Meiryo,system-ui,
    -apple-system,"Segoe UI",Roboto,sans-serif;
  --mono:ui-monospace,"SFMono-Regular",Consolas,"Courier New",monospace;
}
*{box-sizing:border-box;}
html{-webkit-text-size-adjust:100%;}
body{margin:0;background:#eef0f3;color:var(--ink);font-family:var(--sans);
  font-size:16px;line-height:1.85;font-feature-settings:"palt";
  -webkit-font-smoothing:antialiased;-moz-osx-font-smoothing:grayscale;text-rendering:optimizeLegibility;}
img{max-width:100%;height:auto;border-radius:8px;}
::selection{background:var(--accent-soft);}
.sheet{max-width:var(--maxw);margin:32px auto;background:var(--paper);
  padding:56px 64px;border:1px solid var(--line);border-radius:var(--radius);
  box-shadow:0 1px 3px rgba(0,0,0,.05),0 14px 32px rgba(20,30,60,.07);}
.doc-header{margin:0 0 32px;padding-bottom:16px;border-bottom:2px solid var(--accent);}
.doc-title{font-size:28px;font-weight:800;line-height:1.35;margin:0;letter-spacing:.01em;}
.doc-date{color:var(--muted);font-size:13px;margin-top:10px;}
.doc-body>*:first-child{margin-top:0;}
h1,h2,h3,h4,h5,h6{line-height:1.4;font-weight:700;}
h1{font-size:23px;margin:1.8em 0 .6em;padding-bottom:.3em;border-bottom:1px solid var(--line);}
h2{font-size:19.5px;margin:1.7em 0 .5em;padding-left:12px;border-left:4px solid var(--accent);}
h3{font-size:16.5px;margin:1.5em 0 .4em;color:#2b3138;}
h4{font-size:15px;margin:1.3em 0 .3em;color:var(--muted);}
h5{font-size:14px;margin:1.2em 0 .3em;color:var(--muted);}
h6{font-size:13px;margin:1.1em 0 .3em;color:var(--muted);font-weight:600;}
p{margin:.9em 0;}
ul,ol{margin:.7em 0;padding-left:1.6em;}
li{margin:.35em 0;}
li::marker{color:var(--accent);}
li.task{list-style:none;margin-left:-1.5em;}
li.task>input{margin-right:.5em;vertical-align:middle;}
a{color:var(--accent);text-decoration:none;border-bottom:1px solid transparent;}
a:hover{border-bottom-color:currentColor;}
strong{font-weight:700;}
blockquote{margin:1.1em 0;padding:12px 18px;border-left:4px solid var(--accent);
  background:var(--accent-soft);border-radius:0 8px 8px 0;color:#39414b;}
blockquote p{margin:.3em 0;}
hr{border:none;border-top:1px solid var(--line);margin:2em 0;}
:not(pre)>code{font-family:var(--mono);font-size:.86em;background:var(--accent-soft);
  color:#33373d;padding:.15em .42em;border-radius:5px;}
pre{background:var(--soft);border:1px solid var(--line);border-radius:10px;
  padding:14px 16px;overflow-x:auto;margin:1.1em 0;}
pre code{font-family:var(--mono);font-size:13.5px;line-height:1.6;color:#1f2328;}
pre code.hljs{background:transparent;padding:0;}
pre.mermaid{background:none;border:none;padding:0;overflow:visible;text-align:center;
  margin:1.4em 0;line-height:normal;}
.mermaid svg{max-width:100%;height:auto;}
.table-wrap{overflow-x:auto;margin:1.2em 0;}
table{width:100%;border-collapse:separate;border-spacing:0;font-size:14.5px;
  border:1px solid var(--line);border-radius:8px;overflow:hidden;}
th,td{padding:9px 14px;text-align:left;border-bottom:1px solid var(--line);vertical-align:top;}
th{background:var(--soft);font-weight:700;white-space:nowrap;}
tbody tr:nth-child(even){background:var(--zebra);}
tbody tr:last-child td{border-bottom:none;}
.doc-footer{margin-top:40px;padding-top:14px;border-top:1px solid var(--line);
  color:var(--muted);font-size:12px;text-align:right;}
@media (max-width:640px){.sheet{margin:0;border:none;border-radius:0;padding:28px 20px;}}
@media print{
  body{background:#fff;}
  .sheet{margin:0;max-width:none;border:none;border-radius:0;box-shadow:none;padding:0;}
  h1,h2,h3,h4,h5,h6{break-after:avoid;}
  table,pre,blockquote,img{break-inside:avoid;}
  .doc-footer{display:none;}
}
.ref-figs{display:flex;flex-direction:column;gap:18px;margin-top:8px;}
.ref-figs figure{margin:0;}
.ref-figs img{max-width:100%;border:1px solid #e7e9ee;border-radius:8px;}
.ref-figs figcaption{color:var(--muted);font-size:12.5px;margin-top:6px;}
""".strip()


def _decode_figures(figures: list | None, cap: int = 8) -> list[tuple[bytes, str]]:
    """出典の図 [{data(base64), caption}] を [(bytes, caption)] に復元する(不正は無視)。

    Word/PowerPoint/PDF ライブラリが扱えない形式(WebP/BMP 等)は PNG へ変換して
    互換性を確保する(変換できない図は黙って除外し、文書生成は止めない)。
    """
    out: list[tuple[bytes, str]] = []
    for f in (figures or [])[:cap]:
        try:
            data = base64.b64decode(str((f or {}).get("data") or ""))
            if len(data) <= 100:
                continue
            if not (data[:8] == b"\x89PNG\r\n\x1a\n" or data[:3] == b"\xff\xd8\xff"):
                data = _to_png(data)
                if not data:
                    continue
            out.append((data, str((f or {}).get("caption") or "").strip()))
        except Exception:
            log.debug("_decode_figures: 例外を無視して継続", exc_info=True)
    return out


def _to_png(data: bytes) -> bytes:
    """PNG/JPEG 以外の画像バイト列を PNG へ変換する(失敗時は b'')。"""
    try:
        from PIL import Image as PILImage
        img = PILImage.open(io.BytesIO(data))
        if img.mode in ("P", "LA"):
            img = img.convert("RGBA")
        elif img.mode not in ("RGB", "RGBA", "L"):
            img = img.convert("RGB")
        buf = io.BytesIO()
        img.save(buf, format="PNG")
        return buf.getvalue()
    except Exception:
        log.debug("_to_png: 変換に失敗(この図は除外)", exc_info=True)
        return b""


def _img_mime(data: bytes) -> str:
    if data[:8] == b"\x89PNG\r\n\x1a\n":
        return "image/png"
    if data[:6] in (b"GIF87a", b"GIF89a"):
        return "image/gif"
    if data[:4] == b"RIFF" and data[8:12] == b"WEBP":
        return "image/webp"
    return "image/jpeg"


def _figure_dims(data: bytes) -> tuple[float, float]:
    """画像の (幅px, 高さpx)。Pillow が無い・読めない場合は妥当な既定値。"""
    try:
        from PIL import Image as PILImage
        with PILImage.open(io.BytesIO(data)) as img:
            return float(img.size[0]), float(img.size[1])
    except Exception:
        log.debug("_figure_dims: 例外を無視して継続", exc_info=True)
        return 600.0, 400.0


def to_html(md: str, title: str = "回答", figures: list | None = None) -> bytes:
    blocks = parse_blocks(md)
    # 先頭の見出しがタイトルと同じなら重複を避けて省く
    if blocks and blocks[0]["type"] == "heading" and \
            _strip_inline(blocks[0]["text"]).strip() == (title or "").strip():
        blocks = blocks[1:]

    body: list[str] = []
    has_mermaid = False
    has_code = False
    for b in blocks:
        t = b["type"]
        if t == "heading":
            lv = min(b["level"] + 1, 6)   # ページ表題が <h1>。本文は <h2> 起点に降格(h1重複を防ぐ)
            body.append(f"<h{lv}>{_inline_html(b['text'])}</h{lv}>")
        elif t == "paragraph":
            body.append(f"<p>{_inline_html(b['text']).replace(chr(10), '<br>')}</p>")
        elif t == "list":
            body.append(_render_list(b["items"]))
        elif t == "code":
            if b.get("lang", "").lower() == "mermaid":
                # Mermaid 図はコードではなく図として描画する(後段でライブラリを同梱)
                has_mermaid = True
                body.append(f'<pre class="mermaid">{_html.escape(b["text"])}</pre>')
            else:
                has_code = True
                lang = re.sub(r"[^A-Za-z0-9+#-]", "", b.get("lang", "")).lower()
                cls = f' class="language-{lang}"' if lang else ""
                body.append(f"<pre><code{cls}>{_html.escape(b['text'])}</code></pre>")
        elif t == "quote":
            body.append(f"<blockquote>{_inline_html(b['text']).replace(chr(10), '<br>')}</blockquote>")
        elif t == "table":
            al = b.get("aligns", [])
            head = "".join(f"<th{_align_style(al, j)}>{_inline_html(c)}</th>"
                           for j, c in enumerate(b["header"]))
            rows = "".join(
                "<tr>" + "".join(f"<td{_align_style(al, j)}>{_inline_html(c)}</td>"
                                 for j, c in enumerate(r)) + "</tr>"
                for r in b["rows"])
            body.append(f'<div class="table-wrap"><table><thead><tr>{head}</tr></thead>'
                        f"<tbody>{rows}</tbody></table></div>")

    # 参考図(回答の根拠になった文書内の図)。data URI で自己完結させる
    figs = _decode_figures(figures)
    if figs:
        items = []
        for i, (data, cap) in enumerate(figs, 1):
            b64 = base64.b64encode(data).decode("ascii")
            caption = f"図{i}" + (f"　{_html.escape(cap)}" if cap else "")
            items.append(f'<figure><img src="data:{_img_mime(data)};base64,{b64}" '
                         f'alt="参考図{i}"><figcaption>{caption}</figcaption></figure>')
        body.append('<h2>参考図(出典)</h2><div class="ref-figs">' + "".join(items) + "</div>")

    return _html_page(title, "\n".join(body),
                      with_mermaid=has_mermaid, with_code=has_code).encode("utf-8")


# Mermaid 同梱(vendored を1度だけ読み込みキャッシュ)。図を含む HTML のみに埋め込む。
_MERMAID_JS_CACHE: str | None = None


def _mermaid_scripts() -> str:
    """エクスポートHTMLに同梱する Mermaid ライブラリ＋初期化(オフラインでも図が描画される)。"""
    global _MERMAID_JS_CACHE
    if _MERMAID_JS_CACHE is None:
        from pathlib import Path
        p = Path(__file__).resolve().parent.parent / "static" / "vendor" / "mermaid.min.js"
        try:
            js = p.read_text(encoding="utf-8")
            js = re.sub(r"</(script)", r"<\\/\1", js, flags=re.IGNORECASE)  # </script> での早期終了を防ぐ
            _MERMAID_JS_CACHE = js
        except Exception as e:  # 読めなければ図は <pre> のまま(コード表示)になる
            log.warning("mermaid.min.js を同梱できません: %s", e)
            _MERMAID_JS_CACHE = ""
    if not _MERMAID_JS_CACHE:
        return ""
    init = ("mermaid.initialize({startOnLoad:true,securityLevel:'loose',theme:'default',"
            "flowchart:{htmlLabels:true,useMaxWidth:true}});")
    return f"<script>{_MERMAID_JS_CACHE}</script>\n<script>{init}</script>"


# highlight.js 同梱(コードの構文ハイライトをオフラインHTMLに内蔵)
_HLJS_CACHE: tuple[str, str] | None = None


def _highlight_assets() -> tuple[str, str]:
    """(テーマCSS, ライブラリJS) を返す。読めなければ ("", "")。"""
    global _HLJS_CACHE
    if _HLJS_CACHE is None:
        from pathlib import Path
        base = Path(__file__).resolve().parent.parent / "static" / "vendor"
        try:
            js = (base / "highlight.min.js").read_text(encoding="utf-8")
            js = re.sub(r"</(script)", r"<\\/\1", js, flags=re.IGNORECASE)
            css = (base / "github.min.css").read_text(encoding="utf-8")
            _HLJS_CACHE = (css, js)
        except Exception as e:   # 読めなければハイライト無し(コードは素のまま表示)
            log.warning("highlight.js を同梱できません: %s", e)
            _HLJS_CACHE = ("", "")
    return _HLJS_CACHE


def _highlight_head_and_scripts() -> tuple[str, str]:
    """コードを含むHTML用の (head追加CSS, body末尾スクリプト)。"""
    css, js = _highlight_assets()
    if not js:
        return "", ""
    head = f"<style>{css}</style>"
    scripts = (f"<script>{js}</script>\n"
               "<script>hljs.highlightAll();</script>")
    return head, scripts


def _html_page(title: str, body_html: str, with_mermaid: bool = False,
               with_code: bool = False) -> str:
    """整ったドキュメントの外枠(CSS・タイトル・日付・フッター)に本文HTMLを差し込む。"""
    date = datetime.now().strftime("%Y年%m月%d日")
    parts = []
    if with_mermaid:
        parts.append(_mermaid_scripts())
    head_extra = ""
    if with_code:
        hl_head, hl_scripts = _highlight_head_and_scripts()
        head_extra += hl_head
        if hl_scripts:
            parts.append(hl_scripts)
    scripts = "\n".join(p for p in parts if p)
    return f"""<!DOCTYPE html>
<html lang="ja">
<head>
<meta charset="UTF-8">
<meta name="viewport" content="width=device-width, initial-scale=1">
<title>{_html.escape(title)}</title>
<style>{_HTML_CSS}</style>
{head_extra}</head>
<body>
<main class="sheet">
<header class="doc-header">
<h1 class="doc-title">{_html.escape(title)}</h1>
<div class="doc-date">{date}</div>
</header>
<div class="doc-body">
{body_html}
</div>
<footer class="doc-footer">社内文書アシスタントで作成</footer>
</main>
{scripts}
</body>
</html>"""


# モデル生成HTMLが mermaid を使っているかの判定 / 外部・自前ローダの除去
_MERMAID_USE_RE = re.compile(
    r'class\s*=\s*["\']mermaid|<pre[^>]*mermaid|mermaid\.(?:initialize|run|mermaidAPI)'
    r'|import\s+mermaid|["\'][^"\']*mermaid[^"\']*\.(?:m?js)', re.IGNORECASE)


def _uses_mermaid(s: str) -> bool:
    return bool(_MERMAID_USE_RE.search(s or ""))


def _strip_mermaid_loaders(html: str) -> str:
    """外部CDNの mermaid 読み込みと、インラインの import/初期化スクリプトを取り除く。
    (オフラインで動かない/重複初期化するため、同梱版に一本化する)"""
    html = re.sub(r'(?is)<script[^>]+src=["\'][^"\']*mermaid[^"\']*["\'][^>]*>\s*</script>', "", html)
    html = re.sub(r'(?is)<script\b[^>]*>(?:(?!</script>).)*?\bmermaid\b(?:(?!</script>).)*?</script>', "", html)
    return html


def _ensure_mermaid(html: str) -> str:
    """完全HTMLにフローチャート(mermaid)があれば、確実に描画されるよう同梱版を注入する。"""
    if not _uses_mermaid(html):
        return html
    html = _strip_mermaid_loaders(html)
    scripts = _mermaid_scripts()
    if not scripts:
        return html
    m = re.search(r"(?is)</body>", html)   # re.sub は使わない(JS内の \ をテンプレ解釈するため)
    if m:
        return html[:m.start()] + scripts + "\n" + html[m.start():]
    return html + scripts


def _extract_html_document(content: str, title: str) -> str | None:
    """回答に実HTMLが含まれていれば、本物のHTMLページ文字列を返す(無ければ None)。

    - 完全なHTML(<style>あり) → モデルのデザインを尊重して丸ごと出力(図は同梱版で描画)。
    - <body> や HTMLフラグメント → 本文だけ抜き出して当アプリの整ったテンプレに差し込む。
    """
    s = (content or "").strip()
    code = None
    if re.match(r"(?is)^\s*(<!doctype html|<html\b)", s):
        code = s
    else:
        for b in parse_blocks(content):
            if b["type"] == "code":
                bt = b["text"]
                if (b.get("lang", "").lower() in ("html", "htm", "xml")) or \
                        re.search(r"(?is)<!doctype html|<html\b|<body\b|<h[1-6][ >]|<table\b|<div\b", bt):
                    code = bt
                    break
    if not code:
        return None
    has_full = bool(re.search(r"(?is)<html\b", code))
    has_style = bool(re.search(r"(?is)<style\b|<link[^>]+stylesheet", code))
    if has_full and has_style:
        return _ensure_mermaid(code)  # 自前デザインを尊重しつつ、図は同梱版で確実に描画
    m = re.search(r"(?is)<body[^>]*>(.*?)</body>", code)
    if m:
        inner = m.group(1)
    elif not has_full:
        inner = code  # フラグメント
    else:
        inner = re.sub(r"(?is)^.*?<html[^>]*>", "", code)
        inner = re.sub(r"(?is)</html\s*>\s*$", "", inner)
        inner = re.sub(r"(?is)<head\b.*?</head>", "", inner)
    uses = _uses_mermaid(inner)
    inner = _strip_mermaid_loaders(inner)       # フラグメント内のCDN/初期化も同梱版に一本化
    return _html_page(title, inner.strip(), with_mermaid=uses)


def to_docx(md: str, title: str = "回答", images: list | None = None,
            figures: list | None = None) -> bytes:
    from docx import Document
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.shared import Inches, Pt
    doc = Document()
    mm_idx = 0
    for b in parse_blocks(md):
        t = b["type"]
        if t == "heading":
            doc.add_heading(b["text"], level=min(b["level"], 4))
        elif t == "paragraph":
            p = doc.add_paragraph()
            _add_runs(p, b["text"])
        elif t == "list":
            for it in b["items"]:
                style = "List Number" if it["ordered"] else "List Bullet"
                p = doc.add_paragraph(style=style)
                if _item_level(it):
                    try:
                        p.paragraph_format.left_indent = Inches(0.25 * _item_level(it))
                    except Exception:
                        log.debug("to_docx: 例外を無視して継続", exc_info=True)
                        pass
                _add_runs(p, _item_plain(it))
        elif t == "code":
            is_mmd = b.get("lang", "").lower() == "mermaid"
            img = images[mm_idx] if (is_mmd and images and mm_idx < len(images)) else None
            if is_mmd:
                mm_idx += 1
            if img and img.get("data"):
                try:
                    iw = float(img.get("w") or 0) or 600.0
                    doc.add_picture(io.BytesIO(base64.b64decode(img["data"])),
                                    width=Inches(min(6.3, iw / 96.0)))   # 本文幅~6.3in、96dpi換算
                    doc.paragraphs[-1].alignment = WD_ALIGN_PARAGRAPH.CENTER
                    continue
                except Exception:
                    log.debug("to_docx: 例外を無視して継続", exc_info=True)
                    pass
            for ln in b["text"].split("\n"):
                p = doc.add_paragraph()
                run = p.add_run(ln)
                run.font.name = "Consolas"
                run.font.size = Pt(9)
        elif t == "quote":
            p = doc.add_paragraph(b["text"])
            p.style = "Intense Quote" if "Intense Quote" in [s.name for s in doc.styles] else p.style
        elif t == "table":
            cols = max(len(b["header"]), max((len(r) for r in b["rows"]), default=0)) or 1
            table = doc.add_table(rows=1, cols=cols)
            try:
                table.style = "Light Grid Accent 1"
            except Exception:
                log.debug("to_docx: 例外を無視して継続", exc_info=True)
                pass
            for j, c in enumerate(b["header"]):
                if j < cols:
                    cell = table.rows[0].cells[j]
                    cell.text = _strip_inline(c)
                    for r in cell.paragraphs[0].runs:
                        r.bold = True
            for row in b["rows"]:
                cells = table.add_row().cells
                for j, c in enumerate(row):
                    if j < cols:
                        cells[j].text = _strip_inline(c)
    # 参考図(回答の根拠になった文書内の図)を末尾にまとめて掲載
    figs = _decode_figures(figures)
    if figs:
        doc.add_heading("参考図(出典)", level=2)
        for i, (data, cap) in enumerate(figs, 1):
            try:
                w, h = _figure_dims(data)
                doc.add_picture(io.BytesIO(data), width=Inches(min(5.8, w / 96.0)))
                doc.paragraphs[-1].alignment = WD_ALIGN_PARAGRAPH.CENTER
                p = doc.add_paragraph()
                run = p.add_run(f"図{i}" + (f" {cap}" if cap else ""))
                run.font.size = Pt(9)
                p.alignment = WD_ALIGN_PARAGRAPH.CENTER
            except Exception:
                log.debug("to_docx: 参考図の埋め込みに失敗(無視)", exc_info=True)
    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


def _add_runs(paragraph, text: str) -> None:
    from docx.shared import Pt
    for seg, fmt in _inline_runs(text):
        run = paragraph.add_run(seg)
        if fmt.get("bold"):
            run.bold = True
        if fmt.get("italic"):
            run.italic = True
        if fmt.get("code"):
            run.font.name = "Consolas"
            run.font.size = Pt(10)


def to_xlsx(md: str, title: str = "回答", figures: list | None = None) -> bytes:
    from openpyxl import Workbook
    from openpyxl.styles import Font
    blocks = parse_blocks(md)
    tables = [b for b in blocks if b["type"] == "table"]
    wb = Workbook()
    wb.remove(wb.active)

    if tables:
        for ti, tb in enumerate(tables, 1):
            ws = wb.create_sheet(title=f"表{ti}"[:31])
            ws.append([_strip_inline(c) for c in tb["header"]])
            for c in ws[1]:
                c.font = Font(bold=True)
            for row in tb["rows"]:
                ws.append([_strip_inline(c) for c in row])
            _autofit(ws)

    # 本文(テーブル以外のテキスト)
    text_lines = []
    for b in blocks:
        if b["type"] == "heading":
            text_lines.append(_strip_inline(b["text"]))
        elif b["type"] == "paragraph":
            text_lines.extend(_strip_inline(b["text"]).split("\n"))
        elif b["type"] == "list":
            text_lines.extend("・" + "  " * _item_level(it) + _strip_inline(_item_plain(it))
                              for it in b["items"])
        elif b["type"] == "code":
            text_lines.extend(b["text"].split("\n"))
        elif b["type"] == "quote":
            text_lines.extend(_strip_inline(b["text"]).split("\n"))
    if text_lines:
        ws = wb.create_sheet(title="本文")
        for ln in text_lines:
            ws.append([ln])
        _autofit(ws)

    # 参考図(回答の根拠になった文書内の図)を専用シートに貼る
    figs = _decode_figures(figures)
    if figs:
        try:
            from openpyxl.drawing.image import Image as XLImage
            ws = wb.create_sheet(title="参考図")
            ws.column_dimensions["A"].width = 90
            row = 1
            for i, (data, cap) in enumerate(figs, 1):
                ws.cell(row=row, column=1, value=f"図{i}" + (f" {cap}" if cap else ""))
                w, h = _figure_dims(data)
                scale = min(1.0, 640.0 / w, 420.0 / h)
                xi = XLImage(io.BytesIO(data))
                xi.width, xi.height = int(w * scale), int(h * scale)
                ws.add_image(xi, f"A{row + 1}")
                row += int(xi.height / 19) + 3   # 既定の行高(~19px)換算で次の図の位置へ
        except Exception:
            log.debug("to_xlsx: 参考図シートの作成に失敗(無視)", exc_info=True)

    if not wb.sheetnames:
        wb.create_sheet(title="本文")
    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


def _autofit(ws) -> None:
    widths: dict[int, int] = {}
    for row in ws.iter_rows():
        for cell in row:
            if cell.value is not None:
                widths[cell.column] = max(widths.get(cell.column, 8),
                                          min(60, len(str(cell.value)) + 2))
    from openpyxl.utils import get_column_letter
    for col, w in widths.items():
        ws.column_dimensions[get_column_letter(col)].width = w


def to_pptx(md: str, title: str = "回答", images: list | None = None,
            figures: list | None = None) -> bytes:
    from pptx import Presentation
    prs = Presentation()
    blocks = parse_blocks(md)

    # タイトルスライド
    first_heading = next((b["text"] for b in blocks if b["type"] == "heading"), title)
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = first_heading or title

    def new_content_slide(heading_text: str):
        s = prs.slides.add_slide(prs.slide_layouts[1])
        s.shapes.title.text = heading_text
        body = s.placeholders[1].text_frame
        body.clear()
        return body

    body = None
    mm_idx = 0
    for b in blocks:
        if b["type"] == "heading":
            body = new_content_slide(b["text"])
        else:
            if b["type"] == "code":
                is_mmd = b.get("lang", "").lower() == "mermaid"
                img = images[mm_idx] if (is_mmd and images and mm_idx < len(images)) else None
                if is_mmd:
                    mm_idx += 1
                if img and img.get("data"):
                    _add_picture_slide(prs, base64.b64decode(img["data"]), img.get("w"), img.get("h"))
                    body = None       # 図は専用スライド。次の内容は新スライドへ
                    continue
            if body is None:
                body = new_content_slide("内容")
            if b["type"] == "paragraph":
                for ln in _strip_inline(b["text"]).split("\n"):
                    if ln.strip():
                        _add_bullet(body, ln, 0)
            elif b["type"] == "list":
                for it in b["items"]:
                    _add_bullet(body, _strip_inline(_item_plain(it)), 1 + _item_level(it))
            elif b["type"] == "code":
                for ln in b["text"].split("\n"):
                    _add_bullet(body, ln, 0, mono=True)
            elif b["type"] == "quote":
                _add_bullet(body, _strip_inline(b["text"]), 0)
            elif b["type"] == "table":
                head = " | ".join(_strip_inline(c) for c in b["header"])
                _add_bullet(body, head, 0)
                for row in b["rows"]:
                    _add_bullet(body, " | ".join(_strip_inline(c) for c in row), 1)
    # 参考図(回答の根拠になった文書内の図)。1枚=1スライドで末尾に追加
    for i, (data, cap) in enumerate(_decode_figures(figures), 1):
        try:
            w, h = _figure_dims(data)
            _add_picture_slide(prs, data, w, h)
            _set_last_slide_caption(prs, f"図{i}" + (f" {cap}" if cap else ""))
        except Exception:
            log.debug("to_pptx: 参考図スライドの作成に失敗(無視)", exc_info=True)
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _set_last_slide_caption(prs, text: str) -> None:
    """直近の図スライドの下部にキャプションを置く。"""
    from pptx.util import Emu, Pt
    s = prs.slides[-1]
    box = s.shapes.add_textbox(Emu(int(prs.slide_width * 0.06)),
                               Emu(int(prs.slide_height * 0.92)),
                               Emu(int(prs.slide_width * 0.88)),
                               Emu(int(prs.slide_height * 0.07)))
    tf = box.text_frame
    tf.text = text
    for r in tf.paragraphs[0].runs:
        r.font.size = Pt(12)


def _add_picture_slide(prs, png_bytes: bytes, w, h) -> None:
    """図(Mermaid)を1枚のスライドに中央配置で貼る。"""
    layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[5]
    s = prs.slides.add_slide(layout)
    sw, sh = prs.slide_width, prs.slide_height
    emu = 9525                      # 96dpi: 1px = 9525 EMU
    nat_w = (float(w or 0) or 600.0) * emu
    nat_h = (float(h or 0) or 400.0) * emu
    scale = min(sw * 0.88 / nat_w, sh * 0.82 / nat_h)
    dw, dh = nat_w * scale, nat_h * scale
    s.shapes.add_picture(io.BytesIO(png_bytes), int((sw - dw) / 2), int((sh - dh) / 2),
                         int(dw), int(dh))


def _add_bullet(text_frame, text: str, level: int, mono: bool = False) -> None:
    from pptx.util import Pt
    if text_frame.paragraphs and text_frame.paragraphs[0].text == "" and len(text_frame.paragraphs) == 1:
        p = text_frame.paragraphs[0]
    else:
        p = text_frame.add_paragraph()
    p.text = text
    p.level = min(level, 4)
    if mono:
        for r in p.runs:
            r.font.name = "Consolas"
            r.font.size = Pt(12)


def _pdf_inline(text: str) -> str:
    """Markdown のインライン記法を reportlab のミニマークアップに変換(エスケープ込み)。"""
    s = _LINK.sub(r"\1 (\2)", text or "")
    s = _html.escape(s, quote=False)
    s = re.sub(r"\*\*(.+?)\*\*", r"<b>\1</b>", s)
    s = re.sub(r"(?<!\*)\*(?!\*)([^*]+)\*(?!\*)", r"<i>\1</i>", s)
    s = re.sub(r"`([^`]+)`", r"\1", s)        # インラインコードは素のテキスト(JPフォント維持)
    return s


def to_pdf(md: str, title: str = "回答", images: list | None = None,
           figures: list | None = None) -> bytes:
    """Markdown を、HTML出力と同じ意匠の整ったPDFにする(日本語=内蔵CIDフォント)。
    images(フロントで描画した Mermaid 図のPNG。出現順)が渡されれば図を埋め込む。
    figures(出典の文書内画像)が渡されれば「参考図」として末尾に掲載する。
    画像が無い図はコード枠で表示する。"""
    from reportlab.lib import colors
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.styles import ParagraphStyle
    from reportlab.lib.units import mm
    from reportlab.pdfbase import pdfmetrics
    from reportlab.pdfbase.cidfonts import UnicodeCIDFont
    from reportlab.platypus import (HRFlowable, Image as RLImage, ListFlowable, ListItem,
                                    PageBreak, Paragraph, Preformatted, SimpleDocTemplate, Spacer, Table, TableStyle)
    from reportlab.platypus.tableofcontents import TableOfContents

    FONT = "HeiseiKakuGo-W5"
    for f in (FONT, "HeiseiMin-W3"):
        try:
            pdfmetrics.registerFont(UnicodeCIDFont(f))
        except Exception:
            log.debug("to_pdf: 例外を無視して継続", exc_info=True)
            pass

    ACCENT = colors.HexColor("#3b5bdb"); INK = colors.HexColor("#1f2328")
    MUTED = colors.HexColor("#5b6570"); LINE = colors.HexColor("#e7e9ee")
    SOFT = colors.HexColor("#f6f8fb"); ZEBRA = colors.HexColor("#fafbfd")

    body = ParagraphStyle("body", fontName=FONT, fontSize=10.5, leading=17, textColor=INK, spaceAfter=6)
    heads = {
        1: ParagraphStyle("h1", fontName=FONT, fontSize=16, leading=22, textColor=INK, spaceBefore=14, spaceAfter=6),
        2: ParagraphStyle("h2", fontName=FONT, fontSize=14, leading=20, textColor=ACCENT, spaceBefore=12, spaceAfter=5),
        3: ParagraphStyle("h3", fontName=FONT, fontSize=12, leading=18, textColor=INK, spaceBefore=10, spaceAfter=4),
        4: ParagraphStyle("h4", fontName=FONT, fontSize=11, leading=16, textColor=MUTED, spaceBefore=8, spaceAfter=3),
    }
    code_st = ParagraphStyle("code", fontName=FONT, fontSize=9, leading=13.5, textColor=INK,
                             backColor=SOFT, borderColor=LINE, borderWidth=0.5, borderPadding=8, spaceAfter=8)
    quote_st = ParagraphStyle("quote", fontName=FONT, fontSize=10.5, leading=17,
                              textColor=colors.HexColor("#39414b"), leftIndent=10,
                              backColor=colors.HexColor("#eef2ff"), borderPadding=8, spaceAfter=8)
    title_st = ParagraphStyle("title", fontName=FONT, fontSize=20, leading=26, textColor=INK, spaceAfter=2)
    date_st = ParagraphStyle("date", fontName=FONT, fontSize=9, textColor=MUTED, spaceAfter=10)
    toc_head_st = ParagraphStyle("tochead", fontName=FONT, fontSize=14, leading=20,
                                 textColor=INK, spaceBefore=2, spaceAfter=8)

    _outline_levels: dict = {}   # md見出しレベル → 0始まりの階層(blocks 確定後に設定)

    class _PDFDoc(SimpleDocTemplate):
        """見出しを目次(TOCEntry)と PDF しおり(アウトライン)に登録する。"""
        _last_ol = -1

        def beforeDocument(self):       # multiBuild の各パス開始時に階層トラッカをリセット
            self._last_ol = -1

        def afterFlowable(self, flowable):
            if not isinstance(flowable, Paragraph):
                return
            name = flowable.style.name
            if not (len(name) == 2 and name[0] == "h" and name[1].isdigit()):
                return
            ol = _outline_levels.get(int(name[1]), 0)
            ol = min(ol, self._last_ol + 1)   # reportlab: 階層の飛び級(例 0→2)を防ぐ
            self._last_ol = ol
            text = flowable.getPlainText()
            key = "sec-%d" % id(flowable)
            self.canv.bookmarkPage(key)
            self.canv.addOutlineEntry(text, key, level=ol, closed=(ol > 0))
            self.notify("TOCEntry", (min(ol, 2), text, self.page, key))

    def _footer(canvas, doc):           # 全ページにページ番号
        canvas.saveState()
        canvas.setFont(FONT, 8); canvas.setFillColor(MUTED)
        canvas.drawCentredString(A4[0] / 2, 12 * mm, str(canvas.getPageNumber()))
        canvas.restoreState()

    story = [Paragraph(_pdf_inline(title), title_st),
             Paragraph(datetime.now().strftime("%Y年%m月%d日"), date_st),
             HRFlowable(width="100%", thickness=1.2, color=ACCENT, spaceAfter=12)]

    blocks = parse_blocks(md)
    if blocks and blocks[0]["type"] == "heading" and \
            _strip_inline(blocks[0]["text"]).strip() == (title or "").strip():
        blocks = blocks[1:]

    # 実在する見出しレベルを 0始まりへ詰める(しおり/目次の階層を整える)
    _hlv = sorted({b["level"] for b in blocks if b["type"] == "heading"})
    _outline_levels.clear()
    _outline_levels.update({lv: i for i, lv in enumerate(_hlv)})

    # 見出しが3つ以上なら目次(クリックで該当ページへ)を付ける
    if sum(1 for b in blocks if b["type"] == "heading") >= 3:
        toc = TableOfContents()
        toc.levelStyles = [
            ParagraphStyle("toc0", fontName=FONT, fontSize=11, leading=18, textColor=INK),
            ParagraphStyle("toc1", fontName=FONT, fontSize=10, leading=16, leftIndent=14, textColor=MUTED),
            ParagraphStyle("toc2", fontName=FONT, fontSize=9.5, leading=15, leftIndent=28, textColor=MUTED),
        ]
        story += [Paragraph("目次", toc_head_st), toc, PageBreak()]

    content_width = A4[0] - 44 * mm   # 図の最大表示幅(左右マージン22mm)
    mm_idx = 0                        # Mermaid 図の出現番号(images と対応づけ)

    for b in blocks:
        t = b["type"]
        if t == "heading":
            story.append(Paragraph(_pdf_inline(b["text"]), heads[min(b["level"], 4)]))
        elif t == "paragraph":
            story.append(Paragraph(_pdf_inline(b["text"]).replace("\n", "<br/>"), body))
        elif t == "list":
            items = [ListItem(Paragraph(_pdf_inline(_item_plain(it)), body)) for it in b["items"]]
            ordered = bool(b["items"]) and b["items"][0]["ordered"]
            story.append(ListFlowable(items, bulletType="1" if ordered else "bullet",
                                      bulletColor=ACCENT, bulletFontName=FONT, leftIndent=16))
        elif t == "code":
            is_mmd = b.get("lang", "").lower() == "mermaid"
            img = images[mm_idx] if (is_mmd and images and mm_idx < len(images)) else None
            if is_mmd:
                mm_idx += 1
            embedded = False
            if img and img.get("data"):
                try:
                    raw = base64.b64decode(img["data"])
                    iw = float(img.get("w") or 0) or 600.0
                    ih = float(img.get("h") or 0) or 400.0
                    disp_w = min(content_width, iw * 0.75)        # px(96dpi)→pt(72)で実寸、幅で制限
                    story.append(RLImage(io.BytesIO(raw), width=disp_w, height=disp_w * (ih / iw), hAlign="CENTER"))
                    story.append(Spacer(1, 8))
                    embedded = True
                except Exception:
                    log.debug("to_pdf: 例外を無視して継続", exc_info=True)
                    embedded = False
            if not embedded:
                story.append(Preformatted(b["text"], code_st))   # 画像が無い図はコード枠で表示
        elif t == "quote":
            story.append(Paragraph(_pdf_inline(b["text"]).replace("\n", "<br/>"), quote_st))
        elif t == "table":
            data = ([[Paragraph(_pdf_inline(c), body) for c in b["header"]]]
                    + [[Paragraph(_pdf_inline(c), body) for c in r] for r in b["rows"]])
            tbl = Table(data, repeatRows=1, hAlign="LEFT")
            tbl.setStyle(TableStyle([
                ("FONTNAME", (0, 0), (-1, -1), FONT), ("FONTSIZE", (0, 0), (-1, -1), 9.5),
                ("BACKGROUND", (0, 0), (-1, 0), SOFT), ("TEXTCOLOR", (0, 0), (-1, 0), INK),
                ("BOX", (0, 0), (-1, -1), 0.5, LINE), ("LINEBELOW", (0, 0), (-1, -1), 0.4, LINE),
                ("VALIGN", (0, 0), (-1, -1), "TOP"),
                ("LEFTPADDING", (0, 0), (-1, -1), 7), ("RIGHTPADDING", (0, 0), (-1, -1), 7),
                ("TOPPADDING", (0, 0), (-1, -1), 5), ("BOTTOMPADDING", (0, 0), (-1, -1), 5),
                ("ROWBACKGROUNDS", (0, 1), (-1, -1), [colors.white, ZEBRA]),
            ]))
            story.append(tbl)
            story.append(Spacer(1, 6))

    # 参考図(回答の根拠になった文書内の図)を末尾に掲載
    figs = _decode_figures(figures)
    if figs:
        story.append(Paragraph("参考図(出典)", heads[2]))
        cap_st = ParagraphStyle("figcap", fontName=FONT, fontSize=9, leading=13,
                                textColor=MUTED, spaceAfter=10, alignment=1)
        for i, (fdata, cap) in enumerate(figs, 1):
            try:
                iw, ih = _figure_dims(fdata)
                disp_w = min(content_width, iw * 0.75)
                story.append(RLImage(io.BytesIO(fdata), width=disp_w,
                                     height=disp_w * (ih / iw), hAlign="CENTER"))
                story.append(Paragraph(_html.escape(f"図{i}" + (f" {cap}" if cap else "")), cap_st))
            except Exception:
                log.debug("to_pdf: 参考図の埋め込みに失敗(無視)", exc_info=True)

    buf = io.BytesIO()
    doc = _PDFDoc(buf, pagesize=A4, title=title,
                  leftMargin=22 * mm, rightMargin=22 * mm, topMargin=20 * mm, bottomMargin=18 * mm)
    # multiBuild: 目次のページ番号を解決するため複数パスで組版
    doc.multiBuild(story, onFirstPage=_footer, onLaterPages=_footer)
    return buf.getvalue()


def to_csv(md: str, title: str = "回答") -> bytes:
    """表をCSVに(複数表は空行区切り)。表が無ければ本文を1列で出力。Excelで開ける UTF-8 BOM 付き。"""
    import csv
    blocks = parse_blocks(md)
    tables = [b for b in blocks if b["type"] == "table"]
    out = io.StringIO()
    w = csv.writer(out)
    if tables:
        for ti, tb in enumerate(tables):
            if ti > 0:
                w.writerow([])
            w.writerow([_strip_inline(c) for c in tb["header"]])
            for r in tb["rows"]:
                w.writerow([_strip_inline(c) for c in r])
    else:
        for b in blocks:
            if b["type"] in ("heading", "paragraph", "quote"):
                for ln in _strip_inline(b["text"]).split("\n"):
                    w.writerow([ln])
            elif b["type"] == "list":
                for it in b["items"]:
                    w.writerow([_strip_inline(_item_plain(it))])
    return out.getvalue().encode("utf-8-sig")   # BOM付きで Excel の文字化けを防ぐ


# ============================================================
#  プレーンテキスト(.txt)
# ============================================================
def _wcwidth(s: str) -> int:
    """文字列の表示幅(全角=2 / 半角=1)。等幅前提の桁そろえ用。"""
    return sum(2 if unicodedata.east_asian_width(c) in ("W", "F") else 1 for c in s)


def _table_to_txt(header: list[str], rows: list[list[str]]) -> str:
    """表を桁そろえしたプレーンテキストにする(データ欠落を防ぐ。全角幅も考慮)。"""
    table = [header] + rows
    ncol = max((len(r) for r in table), default=0)
    if ncol == 0:
        return ""
    norm = [[_strip_inline(r[j]) if j < len(r) else "" for j in range(ncol)] for r in table]
    widths = [max((_wcwidth(row[j]) for row in norm), default=0) for j in range(ncol)]

    def pad(cell: str, w: int) -> str:
        return cell + " " * max(0, w - _wcwidth(cell))

    def fmt(row: list[str]) -> str:
        return " | ".join(pad(row[j], widths[j]) for j in range(ncol)).rstrip()

    out = [fmt(norm[0]), "-+-".join("-" * w for w in widths)]
    out += [fmt(r) for r in norm[1:]]
    return "\n".join(out)


def _list_to_txt(items: list[dict]) -> str:
    """リストをネスト・番号・タスクを保ったプレーンテキストにする。"""
    lines: list[str] = []
    counters: dict[int, int] = {}
    for it in items:
        lvl = _item_level(it)
        for d in [k for k in counters if k > lvl]:   # 深い階層の番号はリセット
            del counters[d]
        if it.get("ordered"):
            counters[lvl] = counters.get(lvl, 0) + 1
            marker = f"{counters[lvl]}. "
        else:
            counters[lvl] = 0                         # 後続の番号付きが1から始まるように
            marker = "・"
        lines.append("    " * lvl + marker + _strip_inline(_item_plain(it)))
    return "\n".join(lines)


def to_txt(md: str, title: str = "回答") -> bytes:
    """Markdown を読みやすいプレーンテキストへ整形する。

    ブロック間を空行で区切り、見出し・段落・リスト(ネスト/番号/タスク)・表・コード・
    引用をすべて保持する(旧実装は表が欠落し、ブロックが詰まって読みにくかった)。
    """
    parts: list[str] = []
    for b in parse_blocks(md):
        t = b["type"]
        if t in ("heading", "paragraph"):
            parts.append(_strip_inline(b["text"]))
        elif t == "quote":
            parts.append("\n".join("> " + _strip_inline(ln) for ln in b["text"].split("\n")))
        elif t == "code":
            parts.append(b.get("text", ""))
        elif t == "list":
            parts.append(_list_to_txt(b["items"]))
        elif t == "table":
            parts.append(_table_to_txt(b.get("header", []), b.get("rows", [])))
    text = "\n\n".join(p for p in parts if p.strip() != "")
    return (text or (md or "")).encode("utf-8")


# ============================================================
#  ファイル名 / ディスパッチ
# ============================================================
_RESERVED_NAMES = {"CON", "PRN", "AUX", "NUL",
                   *(f"COM{i}" for i in range(1, 10)),
                   *(f"LPT{i}" for i in range(1, 10))}


def safe_stem(title: str, limit: int = 40) -> str:
    """ダウンロードファイル名(拡張子なし)に使える安全な語幹を返す。

    OS の禁止文字・制御文字を除去し、Windows で問題になる末尾のドット/空白・
    予約デバイス名(CON/PRN/NUL 等)も回避する。空になる場合は「回答」を返す。
    """
    s = re.sub(r'[\\/:*?"<>|\x00-\x1f]', "", title or "回答")
    s = s.strip()[:limit].strip().rstrip(" .")
    if not s or s.upper() in _RESERVED_NAMES:
        return "回答"
    return s


def export_content(content: str, fmt: str, ext: str | None = None,
                   title: str = "回答", images: list | None = None,
                   figures: list | None = None) -> tuple[bytes, str, str]:
    """(bytes, mime, 拡張子) を返す。figures は出典の文書内画像(参考図として掲載)。"""
    fmt = (fmt or "md").lower()
    title = (title or "回答").strip() or "回答"

    if fmt == "md":
        return content.encode("utf-8"), MIME["md"], "md"
    if fmt == "txt":
        return to_txt(content, title), MIME["txt"], "txt"
    if fmt == "code":
        safe_ext = re.sub(r"[^A-Za-z0-9]", "", (ext or "txt")) or "txt"
        return content.encode("utf-8"), MIME["code"], safe_ext
    if fmt == "html":
        # 回答に実HTMLが含まれていれば、コードとして見せずに本物のHTMLページとして出力
        doc = _extract_html_document(content, title)
        if doc is not None:
            return doc.encode("utf-8"), MIME["html"], "html"
        return to_html(content, title, figures=figures), MIME["html"], "html"
    if fmt == "pdf":
        return to_pdf(content, title, images=images, figures=figures), MIME["pdf"], "pdf"
    if fmt == "csv":
        return to_csv(content, title), MIME["csv"], "csv"
    if fmt == "docx":
        return to_docx(content, title, images=images, figures=figures), MIME["docx"], "docx"
    if fmt == "xlsx":
        return to_xlsx(content, title, figures=figures), MIME["xlsx"], "xlsx"
    if fmt == "pptx":
        return to_pptx(content, title, images=images, figures=figures), MIME["pptx"], "pptx"

    raise ValueError(f"未対応の形式: {fmt}")
