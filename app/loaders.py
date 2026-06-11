"""
loaders.py
拡張子ごとの文書ローダ。各ローダは
  [{"text": ..., "source": ファイル名, "loc": 場所表記}, ...]
を返す。重い依存は関数内で遅延インポートする。
"""
from __future__ import annotations

from pathlib import Path

from . import ocr
from .config import settings
from .logging_setup import get_logger

log = get_logger("loaders")


def _page_text_with_ocr(page, source_name: str, page_no: int) -> str:
    """1ページのテキストを返す。テキスト層が無い(=スキャン)頁だけ OCR にルーティングする。

    page は PyMuPDF の Page を想定するが、get_text()/get_pixmap() を持つ任意のオブジェクトで
    よい(テスト用に fitz 非依存で差し替え可能)。
    """
    text = page.get_text().strip()
    if ocr.needs_ocr(text):   # OCR無効時は常に False=従来動作
        try:
            png = page.get_pixmap(dpi=settings.ocr_dpi).tobytes("png")
            ocr_text = ocr.ocr_image_png(png)
        except Exception as e:   # 画像化失敗等は従来どおり skip 相当
            log.warning("ページ画像化に失敗 %s p.%d: %s", source_name, page_no, e)
            ocr_text = ""
        if ocr_text:
            text = ocr_text
    return text


def load_pdf(path: Path) -> list[dict]:
    import fitz  # PyMuPDF

    docs = []
    with fitz.open(path) as pdf:
        for i, page in enumerate(pdf):
            # テキスト層がある頁はそのまま、無い頁のみ OCR(OCR無効時は常に従来動作)
            text = _page_text_with_ocr(page, path.name, i + 1)
            if text:
                docs.append({"text": text, "source": path.name, "loc": f"p.{i + 1}"})
    return docs


def load_docx(path: Path) -> list[dict]:
    from docx import Document
    doc = Document(str(path))
    docs: list[dict] = []
    buf: list[str] = []
    head = "本文"
    for para in doc.paragraphs:
        t = para.text.strip()
        if not t:
            continue
        if para.style and para.style.name and para.style.name.startswith(("Heading", "見出し")):
            if buf:
                docs.append({"text": "\n".join(buf), "source": path.name, "loc": head})
                buf = []
            head = t
        buf.append(t)
    if buf:
        docs.append({"text": "\n".join(buf), "source": path.name, "loc": head})
    # 表も拾う
    for ti, table in enumerate(doc.tables, 1):
        rows = []
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells if c.text.strip()]
            if cells:
                rows.append(" | ".join(cells))
        if rows:
            docs.append({"text": "\n".join(rows), "source": path.name, "loc": f"表{ti}"})
    return docs


def load_xlsx(path: Path) -> list[dict]:
    import openpyxl
    wb = openpyxl.load_workbook(path, data_only=True, read_only=True)
    docs = []
    for ws in wb.worksheets:
        rows = []
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) for c in row if c is not None]
            if cells:
                rows.append(" | ".join(cells))
        if rows:
            docs.append({"text": "\n".join(rows), "source": path.name, "loc": f"シート:{ws.title}"})
    wb.close()
    return docs


def load_pptx(path: Path) -> list[dict]:
    from pptx import Presentation
    prs = Presentation(str(path))
    docs = []
    for i, slide in enumerate(prs.slides, 1):
        parts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = "".join(r.text for r in para.runs).strip()
                    if t:
                        parts.append(t)
        if parts:
            docs.append({"text": "\n".join(parts), "source": path.name, "loc": f"スライド{i}"})
    return docs


def load_text(path: Path) -> list[dict]:
    raw = path.read_bytes()
    text = None
    for enc in ("utf-8", "utf-8-sig", "cp932", "shift_jis", "euc_jp", "latin-1"):
        try:
            text = raw.decode(enc)
            break
        except (UnicodeDecodeError, LookupError):
            log.debug("load_text: 例外を無視して継続", exc_info=True)
            continue
    if text is None:
        text = raw.decode("utf-8", errors="replace")
    text = text.strip()
    if not text:
        return []
    return [{"text": text, "source": path.name, "loc": "本文"}]


# 拡張子 -> ローダ
LOADERS = {
    ".pdf": load_pdf,
    ".docx": load_docx,
    ".xlsx": load_xlsx,
    ".pptx": load_pptx,
    ".txt": load_text,
    ".md": load_text,
    ".markdown": load_text,
    ".csv": load_text,
    ".tsv": load_text,
    ".log": load_text,
    ".json": load_text,
}

SUPPORTED_EXTS = set(LOADERS.keys())


def is_temp_artifact(name: str) -> bool:
    """Officeのロック/一時ファイルか(取り込み・件数・一覧から除外する)。

    ~$xxx.xlsx は Excel/Word が開いている間だけ存在する所有者情報ファイルで、
    実体はzipではないため読み込むと必ず失敗する。.~lock.* は LibreOffice の同等品。
    """
    n = (name or "")
    return n.startswith("~$") or n.startswith(".~lock.")


def is_supported(path: Path) -> bool:
    return path.suffix.lower() in SUPPORTED_EXTS and not is_temp_artifact(path.name)


def load_file(path: Path) -> list[dict]:
    """1ファイルを読み込む。未対応・失敗時は空リスト(例外は呼び出し側で扱う)。"""
    loader = LOADERS.get(path.suffix.lower())
    if not loader:
        return []
    return loader(path)
