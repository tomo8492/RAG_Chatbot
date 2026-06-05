"""
export.py
回答(Markdown)を各種ファイル形式へ変換する。
  対応: md / txt / html / docx / xlsx / pptx / code(任意拡張子)
重い依存(python-docx等)は関数内で遅延インポート。
"""
from __future__ import annotations

import base64
import html as _html
import io
import re
from datetime import datetime

from .logging_setup import get_logger

log = get_logger("export")

MIME = {
    "md": "text/markdown; charset=utf-8",
    "txt": "text/plain; charset=utf-8",
    "html": "text/html; charset=utf-8",
    "code": "text/plain; charset=utf-8",
    "pdf": "application/pdf",
    "csv": "text/csv; charset=utf-8",
    "docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    "xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    "pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
}
EXT = {"md": "md", "txt": "txt", "html": "html",
       "docx": "docx", "xlsx": "xlsx", "pptx": "pptx"}


# ============================================================
#  Markdown 簡易ブロック解析
# ============================================================
def parse_blocks(md: str) -> list[dict]:
    lines = (md or "").replace("\r\n", "\n").split("\n")
    blocks: list[dict] = []
    i, n = 0, len(lines)

    def is_table_sep(s: str) -> bool:
        return bool(re.match(r"^\s*\|?[\s:\-\|]+\|?\s*$", s)) and "-" in s

    while i < n:
        line = lines[i]

        # コードフェンス
        m = re.match(r"^\s*```(.*)$", line)
        if m:
            lang = m.group(1).strip()
            i += 1
            buf = []
            while i < n and not re.match(r"^\s*```\s*$", lines[i]):
                buf.append(lines[i]); i += 1
            i += 1  # 終端 ``` をスキップ
            blocks.append({"type": "code", "lang": lang, "text": "\n".join(buf)})
            continue

        # 空行
        if line.strip() == "":
            i += 1
            continue

        # 見出し
        m = re.match(r"^(#{1,6})\s+(.*)$", line)
        if m:
            blocks.append({"type": "heading", "level": len(m.group(1)), "text": m.group(2).strip()})
            i += 1
            continue

        # テーブル
        if "|" in line and i + 1 < n and is_table_sep(lines[i + 1]):
            header = _split_row(line)
            i += 2
            rows = []
            while i < n and "|" in lines[i] and lines[i].strip():
                rows.append(_split_row(lines[i])); i += 1
            blocks.append({"type": "table", "header": header, "rows": rows})
            continue

        # 引用
        if re.match(r"^\s*>\s?", line):
            buf = []
            while i < n and re.match(r"^\s*>\s?", lines[i]):
                buf.append(re.sub(r"^\s*>\s?", "", lines[i])); i += 1
            blocks.append({"type": "quote", "text": "\n".join(buf)})
            continue

        # リスト
        if re.match(r"^\s*([-*+]|\d+[.)])\s+", line):
            items = []
            ordered = bool(re.match(r"^\s*\d+[.)]\s+", line))
            while i < n and re.match(r"^\s*([-*+]|\d+[.)])\s+", lines[i]):
                items.append(re.sub(r"^\s*([-*+]|\d+[.)])\s+", "", lines[i]).strip()); i += 1
            blocks.append({"type": "list", "ordered": ordered, "items": items})
            continue

        # 段落(連続する非空行をまとめる)
        buf = []
        while i < n and lines[i].strip() != "" and not re.match(r"^\s*```", lines[i]) \
                and not re.match(r"^(#{1,6})\s+", lines[i]) \
                and not re.match(r"^\s*([-*+]|\d+[.)])\s+", lines[i]) \
                and not re.match(r"^\s*>\s?", lines[i]):
            buf.append(lines[i]); i += 1
        blocks.append({"type": "paragraph", "text": "\n".join(buf)})

    return blocks


def _split_row(line: str) -> list[str]:
    s = line.strip()
    if s.startswith("|"):
        s = s[1:]
    if s.endswith("|"):
        s = s[:-1]
    return [c.strip() for c in s.split("|")]


# ---- インライン処理 ----
_LINK = re.compile(r"\[([^\]]+)\]\(([^)]+)\)")


def _strip_inline(text: str) -> str:
    """強調記号などを除いたプレーン文字列。"""
    text = _LINK.sub(r"\1 (\2)", text)
    text = re.sub(r"\*\*(.+?)\*\*", r"\1", text)
    text = re.sub(r"`([^`]+)`", r"\1", text)
    text = re.sub(r"(?<!\*)\*(?!\*)([^*]+)\*(?!\*)", r"\1", text)
    text = re.sub(r"__(.+?)__", r"\1", text)
    return text


def _inline_html(text: str) -> str:
    out = _html.escape(text)
    # 角括弧・丸括弧はエスケープされないのでリンク記法を変換可能
    out = re.sub(r"\[([^\]]+)\]\(([^)]+)\)", r'<a href="\2">\1</a>', out)
    out = re.sub(r"\*\*(.+?)\*\*", r"<strong>\1</strong>", out)
    out = re.sub(r"`([^`]+)`", r"<code>\1</code>", out)
    out = re.sub(r"(?<!\*)\*(?!\*)([^*]+)\*(?!\*)", r"<em>\1</em>", out)
    return out


def _inline_runs(text: str) -> list[tuple[str, dict]]:
    """docx 用: (テキスト, {bold,italic,code}) のリスト。"""
    text = _LINK.sub(r"\1 (\2)", text)
    parts: list[tuple[str, dict]] = []
    pattern = re.compile(r"(\*\*.+?\*\*|`[^`]+`|(?<!\*)\*(?!\*)[^*]+\*(?!\*))")
    pos = 0
    for m in pattern.finditer(text):
        if m.start() > pos:
            parts.append((text[pos:m.start()], {}))
        tok = m.group(0)
        if tok.startswith("**"):
            parts.append((tok[2:-2], {"bold": True}))
        elif tok.startswith("`"):
            parts.append((tok[1:-1], {"code": True}))
        else:
            parts.append((tok[1:-1], {"italic": True}))
        pos = m.end()
    if pos < len(text):
        parts.append((text[pos:], {}))
    return parts or [(text, {})]


# ============================================================
#  各形式レンダラ
# ============================================================
_HTML_CSS = """
:root{
  --ink:#1f2328; --muted:#5b6570; --accent:#3b5bdb; --accent-soft:rgba(59,91,219,.08);
  --line:#e7e9ee; --soft:#f6f8fb; --zebra:#fafbfd; --paper:#ffffff;
  --radius:10px; --maxw:820px;
  --sans:"Yu Gothic UI","Hiragino Kaku Gothic ProN","Hiragino Sans",Meiryo,system-ui,
    -apple-system,"Segoe UI",Roboto,sans-serif;
  --mono:ui-monospace,"SFMono-Regular",Consolas,"Courier New",monospace;
}
*{box-sizing:border-box;}
html{-webkit-text-size-adjust:100%;}
body{margin:0;background:#eef0f3;color:var(--ink);font-family:var(--sans);
  font-size:16px;line-height:1.85;font-feature-settings:"palt";
  -webkit-font-smoothing:antialiased;-moz-osx-font-smoothing:grayscale;text-rendering:optimizeLegibility;}
img{max-width:100%;height:auto;border-radius:8px;}
::selection{background:var(--accent-soft);}
.sheet{max-width:var(--maxw);margin:32px auto;background:var(--paper);
  padding:56px 64px;border:1px solid var(--line);border-radius:var(--radius);
  box-shadow:0 1px 3px rgba(0,0,0,.05),0 14px 32px rgba(20,30,60,.07);}
.doc-header{margin:0 0 32px;padding-bottom:16px;border-bottom:2px solid var(--accent);}
.doc-title{font-size:28px;font-weight:800;line-height:1.35;margin:0;letter-spacing:.01em;}
.doc-date{color:var(--muted);font-size:13px;margin-top:10px;}
.doc-body>*:first-child{margin-top:0;}
h1,h2,h3,h4{line-height:1.4;font-weight:700;}
h1{font-size:23px;margin:1.8em 0 .6em;padding-bottom:.3em;border-bottom:1px solid var(--line);}
h2{font-size:19.5px;margin:1.7em 0 .5em;padding-left:12px;border-left:4px solid var(--accent);}
h3{font-size:16.5px;margin:1.5em 0 .4em;color:#2b3138;}
h4{font-size:15px;margin:1.3em 0 .3em;color:var(--muted);}
p{margin:.9em 0;}
ul,ol{margin:.7em 0;padding-left:1.6em;}
li{margin:.35em 0;}
li::marker{color:var(--accent);}
a{color:var(--accent);text-decoration:none;border-bottom:1px solid transparent;}
a:hover{border-bottom-color:currentColor;}
strong{font-weight:700;}
blockquote{margin:1.1em 0;padding:12px 18px;border-left:4px solid var(--accent);
  background:var(--accent-soft);border-radius:0 8px 8px 0;color:#39414b;}
blockquote p{margin:.3em 0;}
hr{border:none;border-top:1px solid var(--line);margin:2em 0;}
:not(pre)>code{font-family:var(--mono);font-size:.86em;background:var(--accent-soft);
  color:#33373d;padding:.15em .42em;border-radius:5px;}
pre{background:var(--soft);border:1px solid var(--line);border-radius:10px;
  padding:14px 16px;overflow-x:auto;margin:1.1em 0;}
pre code{font-family:var(--mono);font-size:13.5px;line-height:1.6;color:#1f2328;}
pre.mermaid{background:none;border:none;padding:0;overflow:visible;text-align:center;
  margin:1.4em 0;line-height:normal;}
.mermaid svg{max-width:100%;height:auto;}
.table-wrap{overflow-x:auto;margin:1.2em 0;}
table{width:100%;border-collapse:separate;border-spacing:0;font-size:14.5px;
  border:1px solid var(--line);border-radius:8px;overflow:hidden;}
th,td{padding:9px 14px;text-align:left;border-bottom:1px solid var(--line);vertical-align:top;}
th{background:var(--soft);font-weight:700;white-space:nowrap;}
tbody tr:nth-child(even){background:var(--zebra);}
tbody tr:last-child td{border-bottom:none;}
.doc-footer{margin-top:40px;padding-top:14px;border-top:1px solid var(--line);
  color:var(--muted);font-size:12px;text-align:right;}
@media (max-width:640px){.sheet{margin:0;border:none;border-radius:0;padding:28px 20px;}}
@media print{
  body{background:#fff;}
  .sheet{margin:0;max-width:none;border:none;border-radius:0;box-shadow:none;padding:0;}
  h1,h2,h3,h4{break-after:avoid;}
  table,pre,blockquote,img{break-inside:avoid;}
  .doc-footer{display:none;}
}
""".strip()


def to_html(md: str, title: str = "回答") -> bytes:
    blocks = parse_blocks(md)
    # 先頭の見出しがタイトルと同じなら重複を避けて省く
    if blocks and blocks[0]["type"] == "heading" and \
            _strip_inline(blocks[0]["text"]).strip() == (title or "").strip():
        blocks = blocks[1:]

    body: list[str] = []
    has_mermaid = False
    for b in blocks:
        t = b["type"]
        if t == "heading":
            lv = min(b["level"], 4)
            body.append(f"<h{lv}>{_inline_html(b['text'])}</h{lv}>")
        elif t == "paragraph":
            body.append(f"<p>{_inline_html(b['text']).replace(chr(10), '<br>')}</p>")
        elif t == "list":
            tag = "ol" if b["ordered"] else "ul"
            items = "".join(f"<li>{_inline_html(it)}</li>" for it in b["items"])
            body.append(f"<{tag}>{items}</{tag}>")
        elif t == "code":
            if b.get("lang", "").lower() == "mermaid":
                # Mermaid 図はコードではなく図として描画する(後段でライブラリを同梱)
                has_mermaid = True
                body.append(f'<pre class="mermaid">{_html.escape(b["text"])}</pre>')
            else:
                body.append(f"<pre><code>{_html.escape(b['text'])}</code></pre>")
        elif t == "quote":
            body.append(f"<blockquote>{_inline_html(b['text']).replace(chr(10), '<br>')}</blockquote>")
        elif t == "table":
            head = "".join(f"<th>{_inline_html(c)}</th>" for c in b["header"])
            rows = "".join("<tr>" + "".join(f"<td>{_inline_html(c)}</td>" for c in r) + "</tr>"
                           for r in b["rows"])
            body.append(f'<div class="table-wrap"><table><thead><tr>{head}</tr></thead>'
                        f"<tbody>{rows}</tbody></table></div>")

    return _html_page(title, "\n".join(body), with_mermaid=has_mermaid).encode("utf-8")


# Mermaid 同梱(vendored を1度だけ読み込みキャッシュ)。図を含む HTML のみに埋め込む。
_MERMAID_JS_CACHE: str | None = None


def _mermaid_scripts() -> str:
    """エクスポートHTMLに同梱する Mermaid ライブラリ＋初期化(オフラインでも図が描画される)。"""
    global _MERMAID_JS_CACHE
    if _MERMAID_JS_CACHE is None:
        from pathlib import Path
        p = Path(__file__).resolve().parent / "static" / "vendor" / "mermaid.min.js"
        try:
            js = p.read_text(encoding="utf-8")
            js = re.sub(r"</(script)", r"<\\/\1", js, flags=re.IGNORECASE)  # </script> での早期終了を防ぐ
            _MERMAID_JS_CACHE = js
        except Exception as e:  # 読めなければ図は <pre> のまま(コード表示)になる
            log.warning("mermaid.min.js を同梱できません: %s", e)
            _MERMAID_JS_CACHE = ""
    if not _MERMAID_JS_CACHE:
        return ""
    init = ("mermaid.initialize({startOnLoad:true,securityLevel:'loose',theme:'default',"
            "flowchart:{htmlLabels:true,useMaxWidth:true}});")
    return f"<script>{_MERMAID_JS_CACHE}</script>\n<script>{init}</script>"


def _html_page(title: str, body_html: str, with_mermaid: bool = False) -> str:
    """整ったドキュメントの外枠(CSS・タイトル・日付・フッター)に本文HTMLを差し込む。"""
    date = datetime.now().strftime("%Y年%m月%d日")
    scripts = _mermaid_scripts() if with_mermaid else ""
    return f"""<!DOCTYPE html>
<html lang="ja">
<head>
<meta charset="UTF-8">
<meta name="viewport" content="width=device-width, initial-scale=1">
<title>{_html.escape(title)}</title>
<style>{_HTML_CSS}</style>
</head>
<body>
<main class="sheet">
<header class="doc-header">
<h1 class="doc-title">{_html.escape(title)}</h1>
<div class="doc-date">{date}</div>
</header>
<div class="doc-body">
{body_html}
</div>
<footer class="doc-footer">社内文書アシスタントで作成</footer>
</main>
{scripts}
</body>
</html>"""


# モデル生成HTMLが mermaid を使っているかの判定 / 外部・自前ローダの除去
_MERMAID_USE_RE = re.compile(
    r'class\s*=\s*["\']mermaid|<pre[^>]*mermaid|mermaid\.(?:initialize|run|mermaidAPI)'
    r'|import\s+mermaid|["\'][^"\']*mermaid[^"\']*\.(?:m?js)', re.IGNORECASE)


def _uses_mermaid(s: str) -> bool:
    return bool(_MERMAID_USE_RE.search(s or ""))


def _strip_mermaid_loaders(html: str) -> str:
    """外部CDNの mermaid 読み込みと、インラインの import/初期化スクリプトを取り除く。
    (オフラインで動かない/重複初期化するため、同梱版に一本化する)"""
    html = re.sub(r'(?is)<script[^>]+src=["\'][^"\']*mermaid[^"\']*["\'][^>]*>\s*</script>', "", html)
    html = re.sub(r'(?is)<script\b[^>]*>(?:(?!</script>).)*?\bmermaid\b(?:(?!</script>).)*?</script>', "", html)
    return html


def _ensure_mermaid(html: str) -> str:
    """完全HTMLにフローチャート(mermaid)があれば、確実に描画されるよう同梱版を注入する。"""
    if not _uses_mermaid(html):
        return html
    html = _strip_mermaid_loaders(html)
    scripts = _mermaid_scripts()
    if not scripts:
        return html
    m = re.search(r"(?is)</body>", html)   # re.sub は使わない(JS内の \ をテンプレ解釈するため)
    if m:
        return html[:m.start()] + scripts + "\n" + html[m.start():]
    return html + scripts


def _extract_html_document(content: str, title: str) -> str | None:
    """回答に実HTMLが含まれていれば、本物のHTMLページ文字列を返す(無ければ None)。

    - 完全なHTML(<style>あり) → モデルのデザインを尊重して丸ごと出力(図は同梱版で描画)。
    - <body> や HTMLフラグメント → 本文だけ抜き出して当アプリの整ったテンプレに差し込む。
    """
    s = (content or "").strip()
    code = None
    if re.match(r"(?is)^\s*(<!doctype html|<html\b)", s):
        code = s
    else:
        for b in parse_blocks(content):
            if b["type"] == "code":
                bt = b["text"]
                if (b.get("lang", "").lower() in ("html", "htm", "xml")) or \
                        re.search(r"(?is)<!doctype html|<html\b|<body\b|<h[1-6][ >]|<table\b|<div\b", bt):
                    code = bt
                    break
    if not code:
        return None
    has_full = bool(re.search(r"(?is)<html\b", code))
    has_style = bool(re.search(r"(?is)<style\b|<link[^>]+stylesheet", code))
    if has_full and has_style:
        return _ensure_mermaid(code)  # 自前デザインを尊重しつつ、図は同梱版で確実に描画
    m = re.search(r"(?is)<body[^>]*>(.*?)</body>", code)
    if m:
        inner = m.group(1)
    elif not has_full:
        inner = code  # フラグメント
    else:
        inner = re.sub(r"(?is)^.*?<html[^>]*>", "", code)
        inner = re.sub(r"(?is)</html\s*>\s*$", "", inner)
        inner = re.sub(r"(?is)<head\b.*?</head>", "", inner)
    uses = _uses_mermaid(inner)
    inner = _strip_mermaid_loaders(inner)       # フラグメント内のCDN/初期化も同梱版に一本化
    return _html_page(title, inner.strip(), with_mermaid=uses)


def to_docx(md: str, title: str = "回答") -> bytes:
    from docx import Document
    from docx.shared import Pt, RGBColor
    doc = Document()
    for b in parse_blocks(md):
        t = b["type"]
        if t == "heading":
            doc.add_heading(b["text"], level=min(b["level"], 4))
        elif t == "paragraph":
            p = doc.add_paragraph()
            _add_runs(p, b["text"])
        elif t == "list":
            style = "List Number" if b["ordered"] else "List Bullet"
            for it in b["items"]:
                p = doc.add_paragraph(style=style)
                _add_runs(p, it)
        elif t == "code":
            for ln in b["text"].split("\n"):
                p = doc.add_paragraph()
                run = p.add_run(ln)
                run.font.name = "Consolas"
                run.font.size = Pt(9)
        elif t == "quote":
            p = doc.add_paragraph(b["text"])
            p.style = "Intense Quote" if "Intense Quote" in [s.name for s in doc.styles] else p.style
        elif t == "table":
            cols = max(len(b["header"]), max((len(r) for r in b["rows"]), default=0)) or 1
            table = doc.add_table(rows=1, cols=cols)
            try:
                table.style = "Light Grid Accent 1"
            except Exception:
                pass
            for j, c in enumerate(b["header"]):
                if j < cols:
                    cell = table.rows[0].cells[j]
                    cell.text = _strip_inline(c)
                    for r in cell.paragraphs[0].runs:
                        r.bold = True
            for row in b["rows"]:
                cells = table.add_row().cells
                for j, c in enumerate(row):
                    if j < cols:
                        cells[j].text = _strip_inline(c)
    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


def _add_runs(paragraph, text: str) -> None:
    from docx.shared import Pt
    for seg, fmt in _inline_runs(text):
        run = paragraph.add_run(seg)
        if fmt.get("bold"):
            run.bold = True
        if fmt.get("italic"):
            run.italic = True
        if fmt.get("code"):
            run.font.name = "Consolas"
            run.font.size = Pt(10)


def to_xlsx(md: str, title: str = "回答") -> bytes:
    from openpyxl import Workbook
    from openpyxl.styles import Font
    blocks = parse_blocks(md)
    tables = [b for b in blocks if b["type"] == "table"]
    wb = Workbook()
    wb.remove(wb.active)

    if tables:
        for ti, tb in enumerate(tables, 1):
            ws = wb.create_sheet(title=f"表{ti}"[:31])
            ws.append([_strip_inline(c) for c in tb["header"]])
            for c in ws[1]:
                c.font = Font(bold=True)
            for row in tb["rows"]:
                ws.append([_strip_inline(c) for c in row])
            _autofit(ws)

    # 本文(テーブル以外のテキスト)
    text_lines = []
    for b in blocks:
        if b["type"] == "heading":
            text_lines.append(_strip_inline(b["text"]))
        elif b["type"] == "paragraph":
            text_lines.extend(_strip_inline(b["text"]).split("\n"))
        elif b["type"] == "list":
            text_lines.extend("・" + _strip_inline(it) for it in b["items"])
        elif b["type"] == "code":
            text_lines.extend(b["text"].split("\n"))
        elif b["type"] == "quote":
            text_lines.extend(_strip_inline(b["text"]).split("\n"))
    if text_lines:
        ws = wb.create_sheet(title="本文")
        for ln in text_lines:
            ws.append([ln])
        _autofit(ws)

    if not wb.sheetnames:
        wb.create_sheet(title="本文")
    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


def _autofit(ws) -> None:
    widths: dict[int, int] = {}
    for row in ws.iter_rows():
        for cell in row:
            if cell.value is not None:
                widths[cell.column] = max(widths.get(cell.column, 8),
                                          min(60, len(str(cell.value)) + 2))
    from openpyxl.utils import get_column_letter
    for col, w in widths.items():
        ws.column_dimensions[get_column_letter(col)].width = w


def to_pptx(md: str, title: str = "回答") -> bytes:
    from pptx import Presentation
    from pptx.util import Pt
    prs = Presentation()
    blocks = parse_blocks(md)

    # タイトルスライド
    first_heading = next((b["text"] for b in blocks if b["type"] == "heading"), title)
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = first_heading or title

    def new_content_slide(heading_text: str):
        s = prs.slides.add_slide(prs.slide_layouts[1])
        s.shapes.title.text = heading_text
        body = s.placeholders[1].text_frame
        body.clear()
        return body

    body = None
    for b in blocks:
        if b["type"] == "heading":
            body = new_content_slide(b["text"])
        else:
            if body is None:
                body = new_content_slide("内容")
            if b["type"] == "paragraph":
                for ln in _strip_inline(b["text"]).split("\n"):
                    if ln.strip():
                        _add_bullet(body, ln, 0)
            elif b["type"] == "list":
                for it in b["items"]:
                    _add_bullet(body, _strip_inline(it), 1)
            elif b["type"] == "code":
                for ln in b["text"].split("\n"):
                    _add_bullet(body, ln, 0, mono=True)
            elif b["type"] == "quote":
                _add_bullet(body, _strip_inline(b["text"]), 0)
            elif b["type"] == "table":
                head = " | ".join(_strip_inline(c) for c in b["header"])
                _add_bullet(body, head, 0)
                for row in b["rows"]:
                    _add_bullet(body, " | ".join(_strip_inline(c) for c in row), 1)
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _add_bullet(text_frame, text: str, level: int, mono: bool = False) -> None:
    from pptx.util import Pt
    if text_frame.paragraphs and text_frame.paragraphs[0].text == "" and len(text_frame.paragraphs) == 1:
        p = text_frame.paragraphs[0]
    else:
        p = text_frame.add_paragraph()
    p.text = text
    p.level = min(level, 4)
    if mono:
        for r in p.runs:
            r.font.name = "Consolas"
            r.font.size = Pt(12)


def _pdf_inline(text: str) -> str:
    """Markdown のインライン記法を reportlab のミニマークアップに変換(エスケープ込み)。"""
    s = _LINK.sub(r"\1 (\2)", text or "")
    s = _html.escape(s, quote=False)
    s = re.sub(r"\*\*(.+?)\*\*", r"<b>\1</b>", s)
    s = re.sub(r"(?<!\*)\*(?!\*)([^*]+)\*(?!\*)", r"<i>\1</i>", s)
    s = re.sub(r"`([^`]+)`", r"\1", s)        # インラインコードは素のテキスト(JPフォント維持)
    return s


def to_pdf(md: str, title: str = "回答", images: list | None = None) -> bytes:
    """Markdown を、HTML出力と同じ意匠の整ったPDFにする(日本語=内蔵CIDフォント)。
    images(フロントで描画した Mermaid 図のPNG。出現順)が渡されれば図を埋め込む。
    画像が無い図はコード枠で表示する。"""
    from reportlab.lib import colors
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.styles import ParagraphStyle
    from reportlab.lib.units import mm
    from reportlab.pdfbase import pdfmetrics
    from reportlab.pdfbase.cidfonts import UnicodeCIDFont
    from reportlab.platypus import (HRFlowable, Image as RLImage, ListFlowable, ListItem,
                                    Paragraph, Preformatted, SimpleDocTemplate, Spacer, Table, TableStyle)

    FONT = "HeiseiKakuGo-W5"
    for f in (FONT, "HeiseiMin-W3"):
        try:
            pdfmetrics.registerFont(UnicodeCIDFont(f))
        except Exception:
            pass

    ACCENT = colors.HexColor("#3b5bdb"); INK = colors.HexColor("#1f2328")
    MUTED = colors.HexColor("#5b6570"); LINE = colors.HexColor("#e7e9ee")
    SOFT = colors.HexColor("#f6f8fb"); ZEBRA = colors.HexColor("#fafbfd")

    body = ParagraphStyle("body", fontName=FONT, fontSize=10.5, leading=17, textColor=INK, spaceAfter=6)
    heads = {
        1: ParagraphStyle("h1", fontName=FONT, fontSize=16, leading=22, textColor=INK, spaceBefore=14, spaceAfter=6),
        2: ParagraphStyle("h2", fontName=FONT, fontSize=14, leading=20, textColor=ACCENT, spaceBefore=12, spaceAfter=5),
        3: ParagraphStyle("h3", fontName=FONT, fontSize=12, leading=18, textColor=INK, spaceBefore=10, spaceAfter=4),
        4: ParagraphStyle("h4", fontName=FONT, fontSize=11, leading=16, textColor=MUTED, spaceBefore=8, spaceAfter=3),
    }
    code_st = ParagraphStyle("code", fontName=FONT, fontSize=9, leading=13.5, textColor=INK,
                             backColor=SOFT, borderColor=LINE, borderWidth=0.5, borderPadding=8, spaceAfter=8)
    quote_st = ParagraphStyle("quote", fontName=FONT, fontSize=10.5, leading=17,
                              textColor=colors.HexColor("#39414b"), leftIndent=10,
                              backColor=colors.HexColor("#eef2ff"), borderPadding=8, spaceAfter=8)
    title_st = ParagraphStyle("title", fontName=FONT, fontSize=20, leading=26, textColor=INK, spaceAfter=2)
    date_st = ParagraphStyle("date", fontName=FONT, fontSize=9, textColor=MUTED, spaceAfter=10)

    story = [Paragraph(_pdf_inline(title), title_st),
             Paragraph(datetime.now().strftime("%Y年%m月%d日"), date_st),
             HRFlowable(width="100%", thickness=1.2, color=ACCENT, spaceAfter=12)]

    blocks = parse_blocks(md)
    if blocks and blocks[0]["type"] == "heading" and \
            _strip_inline(blocks[0]["text"]).strip() == (title or "").strip():
        blocks = blocks[1:]

    content_width = A4[0] - 44 * mm   # 図の最大表示幅(左右マージン22mm)
    mm_idx = 0                        # Mermaid 図の出現番号(images と対応づけ)

    for b in blocks:
        t = b["type"]
        if t == "heading":
            story.append(Paragraph(_pdf_inline(b["text"]), heads[min(b["level"], 4)]))
        elif t == "paragraph":
            story.append(Paragraph(_pdf_inline(b["text"]).replace("\n", "<br/>"), body))
        elif t == "list":
            items = [ListItem(Paragraph(_pdf_inline(it), body)) for it in b["items"]]
            story.append(ListFlowable(items, bulletType="1" if b["ordered"] else "bullet",
                                      bulletColor=ACCENT, bulletFontName=FONT, leftIndent=16))
        elif t == "code":
            is_mmd = b.get("lang", "").lower() == "mermaid"
            img = images[mm_idx] if (is_mmd and images and mm_idx < len(images)) else None
            if is_mmd:
                mm_idx += 1
            embedded = False
            if img and img.get("data"):
                try:
                    raw = base64.b64decode(img["data"])
                    iw = float(img.get("w") or 0) or 600.0
                    ih = float(img.get("h") or 0) or 400.0
                    disp_w = min(content_width, iw * 0.75)        # px(96dpi)→pt(72)で実寸、幅で制限
                    story.append(RLImage(io.BytesIO(raw), width=disp_w, height=disp_w * (ih / iw), hAlign="CENTER"))
                    story.append(Spacer(1, 8))
                    embedded = True
                except Exception:
                    embedded = False
            if not embedded:
                story.append(Preformatted(b["text"], code_st))   # 画像が無い図はコード枠で表示
        elif t == "quote":
            story.append(Paragraph(_pdf_inline(b["text"]).replace("\n", "<br/>"), quote_st))
        elif t == "table":
            data = ([[Paragraph(_pdf_inline(c), body) for c in b["header"]]]
                    + [[Paragraph(_pdf_inline(c), body) for c in r] for r in b["rows"]])
            tbl = Table(data, repeatRows=1, hAlign="LEFT")
            tbl.setStyle(TableStyle([
                ("FONTNAME", (0, 0), (-1, -1), FONT), ("FONTSIZE", (0, 0), (-1, -1), 9.5),
                ("BACKGROUND", (0, 0), (-1, 0), SOFT), ("TEXTCOLOR", (0, 0), (-1, 0), INK),
                ("BOX", (0, 0), (-1, -1), 0.5, LINE), ("LINEBELOW", (0, 0), (-1, -1), 0.4, LINE),
                ("VALIGN", (0, 0), (-1, -1), "TOP"),
                ("LEFTPADDING", (0, 0), (-1, -1), 7), ("RIGHTPADDING", (0, 0), (-1, -1), 7),
                ("TOPPADDING", (0, 0), (-1, -1), 5), ("BOTTOMPADDING", (0, 0), (-1, -1), 5),
                ("ROWBACKGROUNDS", (0, 1), (-1, -1), [colors.white, ZEBRA]),
            ]))
            story.append(tbl)
            story.append(Spacer(1, 6))

    buf = io.BytesIO()
    doc = SimpleDocTemplate(buf, pagesize=A4, title=title,
                            leftMargin=22 * mm, rightMargin=22 * mm, topMargin=20 * mm, bottomMargin=18 * mm)
    doc.build(story)
    return buf.getvalue()


def to_csv(md: str, title: str = "回答") -> bytes:
    """表をCSVに(複数表は空行区切り)。表が無ければ本文を1列で出力。Excelで開ける UTF-8 BOM 付き。"""
    import csv
    blocks = parse_blocks(md)
    tables = [b for b in blocks if b["type"] == "table"]
    out = io.StringIO()
    w = csv.writer(out)
    if tables:
        for ti, tb in enumerate(tables):
            if ti > 0:
                w.writerow([])
            w.writerow([_strip_inline(c) for c in tb["header"]])
            for r in tb["rows"]:
                w.writerow([_strip_inline(c) for c in r])
    else:
        for b in blocks:
            if b["type"] in ("heading", "paragraph", "quote"):
                for ln in _strip_inline(b["text"]).split("\n"):
                    w.writerow([ln])
            elif b["type"] == "list":
                for it in b["items"]:
                    w.writerow([_strip_inline(it)])
    return out.getvalue().encode("utf-8-sig")   # BOM付きで Excel の文字化けを防ぐ


# ============================================================
#  ディスパッチ
# ============================================================
def export_content(content: str, fmt: str, ext: str | None = None,
                   title: str = "回答", images: list | None = None) -> tuple[bytes, str, str]:
    """(bytes, mime, 拡張子) を返す。"""
    fmt = (fmt or "md").lower()
    title = (title or "回答").strip() or "回答"

    if fmt == "md":
        return content.encode("utf-8"), MIME["md"], "md"
    if fmt == "txt":
        # 記法を軽く落としたプレーンテキスト
        plain = "\n".join(_strip_inline(b["text"]) if b["type"] in ("paragraph", "heading", "quote")
                          else ("\n".join("・" + _strip_inline(x) for x in b["items"]) if b["type"] == "list"
                                else b.get("text", "")) for b in parse_blocks(content))
        return (plain or content).encode("utf-8"), MIME["txt"], "txt"
    if fmt == "code":
        safe_ext = re.sub(r"[^A-Za-z0-9]", "", (ext or "txt")) or "txt"
        return content.encode("utf-8"), MIME["code"], safe_ext
    if fmt == "html":
        # 回答に実HTMLが含まれていれば、コードとして見せずに本物のHTMLページとして出力
        doc = _extract_html_document(content, title)
        if doc is not None:
            return doc.encode("utf-8"), MIME["html"], "html"
        return to_html(content, title), MIME["html"], "html"
    if fmt == "pdf":
        return to_pdf(content, title, images=images), MIME["pdf"], "pdf"
    if fmt == "csv":
        return to_csv(content, title), MIME["csv"], "csv"
    if fmt == "docx":
        return to_docx(content, title), MIME["docx"], "docx"
    if fmt == "xlsx":
        return to_xlsx(content, title), MIME["xlsx"], "xlsx"
    if fmt == "pptx":
        return to_pptx(content, title), MIME["pptx"], "pptx"

    raise ValueError(f"未対応の形式: {fmt}")
