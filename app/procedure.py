"""procedure.py
作業手順書を「工程ごと(文章+画像)」の構造に分解する(手順ビューア用)。

  - Excel: 工程番号の列(1,2,3… / ①② / 手順N / (1) 等)を自動検出し、
           行範囲で工程を区切る。画像はアンカー行で該当工程に割り当てる。
           検出できないシートは「シート全体」として返す(フォールバック)。
  - PowerPoint: 1スライド=1工程
  - PDF: 1ページ=1工程
  - Word: 見出し節=1工程

画像は doc_images の正規化・保存(内容ハッシュ)を再利用するため、
チャットの図表示と同じURL(/api/doc-images/…)で配信される。
"""
from __future__ import annotations

import re
from pathlib import Path

from . import doc_images
from .logging_setup import get_logger

log = get_logger("procedure")

MAX_STEPS = 300          # 1シート/文書あたりの最大工程数(暴走防止)
MAX_TEXT = 4000          # 1工程の本文上限(文字)
MAX_IMGS_PER_STEP = 10   # 1工程に表示する最大画像数
_MARK_SCAN_COLS = 4      # 工程番号を探す先頭からの列数

_CIRCLED = {chr(0x2460 + i): i + 1 for i in range(20)}   # ①〜⑳
_ZEN2HAN = str.maketrans("０１２３４５６７８９", "0123456789")
_STEP_RE = re.compile(r"(?:工程|手順|STEP|Step|No\.?)?\s*[((]?(\d{1,3})[)).．.、]?")


def _step_no(cell: str) -> int | None:
    """セルが工程番号(1 / ① / 手順2 / (3) 等)なら番号を、違えば None。"""
    s = (cell or "").strip()
    if not s or len(s) > 8:
        return None
    if s in _CIRCLED:
        return _CIRCLED[s]
    m = _STEP_RE.fullmatch(s.translate(_ZEN2HAN))
    if m:
        n = int(m.group(1))
        return n if 1 <= n <= 500 else None
    return None


def _save_img(iid: str, data: bytes) -> str | None:
    """画像を正規化して保存し、配信用URLを返す(アイコン等は None)。"""
    norm = doc_images._normalize(data)
    if not norm:
        return None
    return "/api/doc-images/" + doc_images._save(iid, norm[0], norm[1])


# ------------------------------------------------------------------
#  Excel(本命: 工程番号列の自動検出 + 行アンカー対応づけ)
# ------------------------------------------------------------------
def _xlsx_sheets(path: Path):
    """[(シート名, rows=[(行idx, 位置つきセル, 連結テキスト)], imgs=[(行idx, bytes)])]"""
    import openpyxl
    wb = openpyxl.load_workbook(path, data_only=True)
    out = []
    for ws in wb.worksheets:
        rows = []
        for r_i, row in enumerate(ws.iter_rows(values_only=True)):
            vals = ["" if c is None else str(c).strip() for c in row]
            texts = [v for v in vals if v]
            if texts:
                rows.append((r_i, vals, " | ".join(texts)))
        imgs = []
        for im in getattr(ws, "_images", []) or []:
            try:
                frm = getattr(im.anchor, "_from", None)
                row0 = int(getattr(frm, "row", 0) or 0)
                ref = getattr(im, "ref", None)
                if ref is not None and hasattr(ref, "read"):
                    ref.seek(0)
                    data = ref.read()
                elif isinstance(ref, (bytes, bytearray)):
                    data = bytes(ref)
                else:
                    data = im._data()
                if data:
                    imgs.append((row0, data))
            except Exception:
                log.debug("_xlsx_sheets: 画像1件の取得に失敗(無視)", exc_info=True)
        out.append((ws.title, rows, imgs))
    wb.close()
    return out


def _detect_marker_col(rows) -> tuple[int | None, list[tuple[int, int]]]:
    """工程番号の列を推定する。戻り値: (列idx, [(行idx, 番号)])。見つからなければ (None, [])。

    先頭数列を走査し、「番号セルが2個以上」かつ「概ね昇順(6割以上)」の列のうち
    最も番号が多い列を採用する。
    """
    best: tuple[int, list] | None = None
    for col in range(_MARK_SCAN_COLS):
        marks = []
        for r_i, vals, _ in rows:
            if col < len(vals):
                n = _step_no(vals[col])
                if n is not None:
                    marks.append((r_i, n))
        if len(marks) < 2:
            continue
        asc = sum(1 for a, b in zip(marks, marks[1:]) if b[1] >= a[1])
        if asc / (len(marks) - 1) < 0.6:
            continue                       # 番号がバラバラ → 工程列ではない(数値データ列等)
        if best is None or len(marks) > len(best[1]):
            best = (col, marks)
    return (best[0], best[1]) if best else (None, [])


def _xlsx_view(iid: str, path: Path) -> list[dict]:
    sheets_out = []
    for title, rows, imgs in _xlsx_sheets(path):
        if not rows and not imgs:
            continue
        col, marks = _detect_marker_col(rows)
        if col is None:
            # フォールバック: シート全体を1ブロックで(画像はアンカー行順)
            text = "\n".join(t for _, _, t in rows)[:MAX_TEXT]
            urls = []
            for _, data in sorted(imgs, key=lambda x: x[0])[:MAX_IMGS_PER_STEP]:
                u = _save_img(iid, data)
                if u:
                    urls.append(u)
            sheets_out.append({"name": title, "fallback": True,
                               "steps": [{"no": "", "title": "シート全体", "text": text,
                                          "images": urls}]})
            continue
        # 工程の行範囲: marker行 〜 次のmarker行の直前。最初のmarkerより上は「前置き」
        bounds = [r for r, _ in marks][:MAX_STEPS]
        steps: list[dict] = []
        intro_rows = [(r, v, t) for (r, v, t) in rows if r < bounds[0]]
        if intro_rows:
            steps.append({"no": "", "title": "前置き",
                          "text": "\n".join(t for _, _, t in intro_rows)[:MAX_TEXT], "images": []})
        for k, (start_r, no) in enumerate(marks[:MAX_STEPS]):
            end_r = marks[k + 1][0] if k + 1 < len(marks) else 10 ** 9
            seg = [(r, v, t) for (r, v, t) in rows if start_r <= r < end_r]
            # 見出し: marker行の番号セル以外のテキスト
            head = ""
            for r, vals, _ in seg:
                if r == start_r:
                    others = [v for i, v in enumerate(vals) if v and i != col]
                    head = " ".join(others)[:80]
                    break
            steps.append({"no": str(no), "title": head,
                          "text": "\n".join(t for _, _, t in seg)[:MAX_TEXT], "images": []})
        # 画像をアンカー行で該当工程へ割り当て(前置き工程があれば最初のmarkerより上はそこへ)
        for row0, data in sorted(imgs, key=lambda x: x[0]):
            u = _save_img(iid, data)
            if not u:
                continue
            tgt = None
            for k in range(len(marks) - 1, -1, -1):
                if row0 >= marks[k][0]:
                    tgt = steps[k + (1 if intro_rows else 0)]
                    break
            if tgt is None:
                tgt = steps[0]
            if len(tgt["images"]) < MAX_IMGS_PER_STEP and u not in tgt["images"]:
                tgt["images"].append(u)
        sheets_out.append({"name": title, "fallback": False, "steps": steps})
    return sheets_out


# ------------------------------------------------------------------
#  PowerPoint / PDF / Word(自然な単位=スライド・ページ・見出し節)
# ------------------------------------------------------------------
def _pptx_view(iid: str, path: Path) -> list[dict]:
    from pptx import Presentation
    prs = Presentation(str(path))
    steps = []
    for i, slide in enumerate(prs.slides, 1):
        texts: list[str] = []
        urls: list[str] = []
        title = ""
        for shape in slide.shapes:
            try:
                if shape.has_text_frame:
                    t = "\n".join(p.text for p in shape.text_frame.paragraphs if p.text.strip())
                    if t:
                        if not title and shape == slide.shapes.title:
                            title = t.splitlines()[0][:80]
                        texts.append(t)
            except Exception:
                pass
            try:
                u = _save_img(iid, shape.image.blob)
                if u and len(urls) < MAX_IMGS_PER_STEP:
                    urls.append(u)
            except Exception:
                continue
        steps.append({"no": str(i), "title": title or f"スライド{i}",
                      "text": "\n".join(texts)[:MAX_TEXT], "images": urls})
        if len(steps) >= MAX_STEPS:
            break
    return [{"name": "スライド", "fallback": False, "steps": steps}]


def _pdf_view(iid: str, path: Path) -> list[dict]:
    import fitz
    steps = []
    with fitz.open(path) as pdf:
        for i, page in enumerate(pdf, 1):
            urls: list[str] = []
            for info in page.get_images(full=True):
                try:
                    ex = pdf.extract_image(info[0])
                    u = _save_img(iid, ex["image"])
                    if u and u not in urls and len(urls) < MAX_IMGS_PER_STEP:
                        urls.append(u)
                except Exception:
                    continue
            steps.append({"no": str(i), "title": f"ページ {i}",
                          "text": (page.get_text() or "").strip()[:MAX_TEXT], "images": urls})
            if len(steps) >= MAX_STEPS:
                break
    return [{"name": "ページ", "fallback": False, "steps": steps}]


def _docx_view(iid: str, path: Path) -> list[dict]:
    from docx import Document
    from docx.oxml.ns import qn
    from docx.text.paragraph import Paragraph
    doc = Document(str(path))
    steps: list[dict] = []

    def new_step(title: str):
        steps.append({"no": str(len(steps) + 1), "title": title[:80], "text": "", "images": []})

    new_step("本文")
    for child in doc.element.body.iterchildren():
        blips = []
        if child.tag == qn("w:p"):
            para = Paragraph(child, doc)
            t = para.text.strip()
            if t and para.style and para.style.name and \
                    para.style.name.startswith(("Heading", "見出し")):
                if len(steps) >= MAX_STEPS:
                    break
                new_step(t)
            elif t:
                steps[-1]["text"] = (steps[-1]["text"] + "\n" + t)[:MAX_TEXT]
            blips = child.xpath(".//a:blip")
        elif child.tag == qn("w:tbl"):
            blips = child.xpath(".//a:blip")
        for blip in blips:
            rid = blip.get(qn("r:embed"))
            if not rid:
                continue
            try:
                u = _save_img(iid, doc.part.related_parts[rid].blob)
                if u and len(steps[-1]["images"]) < MAX_IMGS_PER_STEP:
                    steps[-1]["images"].append(u)
            except Exception:
                continue
    steps = [s for s in steps if s["text"].strip() or s["images"]]
    return [{"name": "本文", "fallback": False, "steps": steps}]


# ------------------------------------------------------------------
#  公開API
# ------------------------------------------------------------------
def build_view(iid: str, path: Path) -> dict:
    """1ファイルを工程ビュー構造に分解する。未対応形式は ValueError。"""
    ext = path.suffix.lower()
    if ext == ".xlsx":
        sheets = _xlsx_view(iid, path)
    elif ext == ".pptx":
        sheets = _pptx_view(iid, path)
    elif ext == ".pdf":
        sheets = _pdf_view(iid, path)
    elif ext == ".docx":
        sheets = _docx_view(iid, path)
    else:
        raise ValueError(f"手順ビューア未対応の形式です: {ext}(対応: xlsx/pptx/pdf/docx)")
    return {"source": path.name, "sheets": sheets}
