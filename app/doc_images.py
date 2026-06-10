"""doc_images.py
文書(Excel / Word / PowerPoint / PDF)内の埋め込み画像を抽出して保存し、
チャンクの出典(シート/ページ/節)に紐付ける。

「手順書に『図1参照』とあるのに図が見えない」問題への対応基盤:
  - インデックス作成時に extract_for_file() で画像を data/doc_images/<index_id>/ へ保存
  - 戻り値の loc→画像ID 対応を、同じ場所のチャンクのメタデータに載せる(rag.build_index)
  - 回答時は出典パネルにサムネイル表示(/api/doc-images/ で配信)

ロゴ・罫線などのノイズは寸法/サイズ/出現頻度で除外する。抽出失敗は本文処理に
影響させない(すべて握りつぶしてログのみ)。
"""
from __future__ import annotations

import hashlib
import io
import shutil
import zipfile
from pathlib import Path

from .config import settings
from .logging_setup import get_logger

log = get_logger("doc_images")

# ノイズ除外と暴走防止の調整値
MIN_DIM = 64            # これ未満の幅/高さはアイコン・罫線とみなし除外(px)
MIN_BYTES = 4096        # これ未満のファイルは除外
MAX_DIM = 1600          # これを超える画像は縮小して保存(表示・VLM入力の双方に十分)
MAX_PER_FILE = 40       # 1ファイルから取り込む最大枚数
MAX_PER_LOC = 8         # 1ページ/シート/節あたりの最大枚数
MAX_REPEAT_LOCS = 3     # 同一画像が多くの場所に現れる(=ロゴ等)場合に記録する場所数の上限

_WEB_EXTS = {"png", "jpg", "jpeg", "gif", "webp"}


def _img_dir(iid: str) -> Path:
    return settings.data_dir / "doc_images" / iid


def _sniff_ext(data: bytes) -> str:
    """先頭バイトから画像形式を推定する(不明は '')。"""
    if data[:8] == b"\x89PNG\r\n\x1a\n":
        return "png"
    if data[:3] == b"\xff\xd8\xff":
        return "jpg"
    if data[:6] in (b"GIF87a", b"GIF89a"):
        return "gif"
    if data[:4] == b"RIFF" and data[8:12] == b"WEBP":
        return "webp"
    if data[:2] == b"BM":
        return "bmp"
    if data[:4] in (b"II*\x00", b"MM\x00*"):
        return "tiff"
    if len(data) > 44 and data[40:44] == b" EMF":
        return "emf"
    if data[:4] == b"\xd7\xcd\xc6\x9a":
        return "wmf"
    return ""


def _normalize(data: bytes) -> tuple[bytes, str] | None:
    """画像を検査し、(保存するバイト列, 拡張子) を返す。除外対象は None。

    - 小さすぎる画像(アイコン・罫線)は除外
    - Webで表示できない形式(BMP/TIFF/EMF等)は PNG へ変換(変換できなければ除外)
    - 大きすぎる画像は MAX_DIM に縮小
    Pillow が無い環境では PNG/JPEG のみ寸法チェックなしで通す(縮退)。
    """
    if len(data) < MIN_BYTES:
        return None
    ext = _sniff_ext(data)
    if not ext:
        return None
    try:
        from PIL import Image as PILImage
    except ImportError:
        log.debug("_normalize: Pillow なし(PNG/JPEGのみ通す)")
        return (data, ext) if ext in ("png", "jpg") else None
    try:
        # 注釈: open() は ImageFile を返すが、convert() の戻り(Image)を再代入するため広い型で受ける
        img: PILImage.Image = PILImage.open(io.BytesIO(data))
        w, h = img.size
        if w < MIN_DIM or h < MIN_DIM:
            return None
        if ext in _WEB_EXTS and max(w, h) <= MAX_DIM:
            return data, ("jpg" if ext == "jpeg" else ext)
        # 変換 or 縮小が必要
        img = img.convert("RGBA") if img.mode in ("P", "LA") else img
        if max(w, h) > MAX_DIM:
            img.thumbnail((MAX_DIM, MAX_DIM))
        buf = io.BytesIO()
        if img.mode == "RGBA":
            img.save(buf, format="PNG")
            return buf.getvalue(), "png"
        img.convert("RGB").save(buf, format="JPEG", quality=88)
        return buf.getvalue(), "jpg"
    except Exception:
        log.debug("_normalize: 画像の検査/変換に失敗(除外)", exc_info=True)
        return None


def _save(iid: str, data: bytes, ext: str) -> str:
    """画像を内容ハッシュ名で保存し、画像ID(<iid>/<name>)を返す。同一内容は再利用。"""
    name = hashlib.sha1(data).hexdigest()[:16] + "." + ext
    d = _img_dir(iid)
    d.mkdir(parents=True, exist_ok=True)
    p = d / name
    if not p.exists():
        p.write_bytes(data)
    return f"{iid}/{name}"


# ------------------------------------------------------------------
#  形式別の抽出((loc, バイト列) のリストを返す。loc は loaders の loc と同じ表記)
# ------------------------------------------------------------------
def _extract_xlsx(path: Path) -> list[tuple[str, bytes]]:
    out: list[tuple[str, bytes]] = []
    try:
        import openpyxl
        wb = openpyxl.load_workbook(path, data_only=True)   # 画像は read_only では取れない
        for ws in wb.worksheets:
            loc = f"シート:{ws.title}"
            for img in getattr(ws, "_images", []) or []:
                try:
                    ref = getattr(img, "ref", None)
                    if ref is not None and hasattr(ref, "read"):
                        ref.seek(0)
                        data = ref.read()
                    elif isinstance(ref, (bytes, bytearray)):
                        data = bytes(ref)
                    else:
                        data = img._data()   # openpyxl の内部ローダ
                    if data:
                        out.append((loc, data))
                except Exception:
                    log.debug("_extract_xlsx: 画像1件の取得に失敗(無視)", exc_info=True)
        wb.close()
    except Exception:
        log.debug("_extract_xlsx: 失敗(zip直読みへフォールバック)", exc_info=True)
        # フォールバック: zip の xl/media を場所情報なしで拾う(リンクは付かないが図検索は可能)
        try:
            with zipfile.ZipFile(path) as z:
                for n in z.namelist():
                    if n.startswith("xl/media/"):
                        out.append(("", z.read(n)))
        except Exception:
            log.debug("_extract_xlsx: zipフォールバックも失敗", exc_info=True)
    return out


def _extract_docx(path: Path) -> list[tuple[str, bytes]]:
    """Word: 段落を順に歩き、現在の見出し(loaders.load_docx と同じ規則)に画像を紐付ける。"""
    out: list[tuple[str, bytes]] = []
    try:
        from docx import Document
        from docx.oxml.ns import qn
        doc = Document(str(path))
        head = "本文"
        for para in doc.paragraphs:
            t = para.text.strip()
            if t and para.style and para.style.name and \
                    para.style.name.startswith(("Heading", "見出し")):
                head = t
            for blip in para._element.xpath(".//a:blip"):
                rid = blip.get(qn("r:embed"))
                if not rid:
                    continue
                try:
                    out.append((head, doc.part.related_parts[rid].blob))
                except Exception:
                    log.debug("_extract_docx: 画像1件の取得に失敗(無視)", exc_info=True)
    except Exception:
        log.debug("_extract_docx: 失敗(画像なしで継続)", exc_info=True)
    return out


def _extract_pptx(path: Path) -> list[tuple[str, bytes]]:
    out: list[tuple[str, bytes]] = []
    try:
        from pptx import Presentation
        prs = Presentation(str(path))
        for i, slide in enumerate(prs.slides, 1):
            for shape in slide.shapes:
                try:
                    img = shape.image          # 画像でない図形は例外になる
                except Exception:
                    continue
                out.append((f"スライド{i}", img.blob))
    except Exception:
        log.debug("_extract_pptx: 失敗(画像なしで継続)", exc_info=True)
    return out


def _extract_pdf(path: Path) -> list[tuple[str, bytes]]:
    out: list[tuple[str, bytes]] = []
    try:
        import fitz
        with fitz.open(path) as pdf:
            # 同じ画像(xref)が多くのページに出る=ヘッダロゴ等 → 除外するため出現数を数える
            pages_of: dict[int, int] = {}
            per_page: list[list[int]] = []
            for page in pdf:
                xrefs = []
                for info in page.get_images(full=True):
                    xref = info[0]
                    xrefs.append(xref)
                    pages_of[xref] = pages_of.get(xref, 0) + 1
                per_page.append(xrefs)
            done: set[int] = set()
            for i, xrefs in enumerate(per_page, 1):
                n = 0
                for xref in xrefs:
                    if xref in done or pages_of.get(xref, 0) > MAX_REPEAT_LOCS:
                        continue
                    try:
                        ex = pdf.extract_image(xref)
                        out.append((f"p.{i}", ex["image"]))
                        done.add(xref)
                        n += 1
                        if n >= MAX_PER_LOC:
                            break
                    except Exception:
                        log.debug("_extract_pdf: 画像1件の取得に失敗(無視)", exc_info=True)
    except Exception:
        log.debug("_extract_pdf: 失敗(画像なしで継続)", exc_info=True)
    return out


_EXTRACTORS = {
    ".xlsx": _extract_xlsx,
    ".docx": _extract_docx,
    ".pptx": _extract_pptx,
    ".pdf": _extract_pdf,
}


def extract_for_file(iid: str, path: Path) -> tuple[dict[str, list[str]], list[tuple[str, bytes, str]]]:
    """1ファイルの埋め込み画像を抽出・保存する。

    戻り値:
      loc_map: {loc(シート/ページ/節): [画像ID, ...]} — チャンクのメタデータ用
      images:  [(画像ID, 保存したバイト列, loc), ...]   — 図の説明文生成(VLM)用
    """
    fn = _EXTRACTORS.get(path.suffix.lower())
    if not fn:
        return {}, []
    raw = fn(path)
    loc_map: dict[str, list[str]] = {}
    images: list[tuple[str, bytes, str]] = []
    seen_locs: dict[str, list[str]] = {}   # 画像ID -> 出現loc(ロゴの繰り返し検出)
    for loc, data in raw:
        if len(images) >= MAX_PER_FILE:
            break
        norm = _normalize(data)
        if not norm:
            continue
        ndata, ext = norm
        img_id = _save(iid, ndata, ext)
        locs = seen_locs.setdefault(img_id, [])
        if loc in locs:
            continue                          # 同じ場所の重複
        if len(locs) >= MAX_REPEAT_LOCS:
            continue                          # 多数の場所に出る=ロゴ等
        locs.append(loc)
        if loc:
            ids = loc_map.setdefault(loc, [])
            if len(ids) < MAX_PER_LOC and img_id not in ids:
                ids.append(img_id)
        if not any(i[0] == img_id for i in images):
            images.append((img_id, ndata, loc))
    return loc_map, images


def delete_index_images(iid: str) -> None:
    """インデックス削除時に画像フォルダも消す。"""
    try:
        shutil.rmtree(_img_dir(iid), ignore_errors=True)
    except Exception:
        log.debug("delete_index_images: 例外を無視して継続", exc_info=True)
